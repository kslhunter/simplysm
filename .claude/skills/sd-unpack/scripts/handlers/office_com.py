"""Office (docx/pptx/xlsx) + 레거시 (doc/ppt/xls/xlsb) COM 핸들러.

시각 산출물은 PNG, 텍스트 산출물은 MD 로 분리:
- docx → pages/<NNN>.png + pages/<NNN>.md (페이지별)
- pptx → slides/<NN>_<title>.png + .md + .notes.md (슬라이드별)
- xlsx → sheets/<NN>_<name>.png + .md + .formulas.json (시트별)

Office COM 호출은 office_worker.py subprocess 로 격리 (cleanup race 회피).
이 모듈 (office_com.py) 은 호출자 + Office 외 작업 (.md, ZIP strip, 매크로 추출, README 생성).
원칙: 처리 실패는 묻지 않고 그대로 throw. try/finally 는 락/임시 폴더 cleanup 에만 사용.
"""
from __future__ import annotations

import json
import os
import sys
import zipfile
from pathlib import Path
from typing import Optional

from . import _common
from .dispatch import maybe_recurse_attachment


# ====================================================================
# 진입점
# ====================================================================

def run(input_path: Path, out_dir: Path) -> None:
    ext = input_path.suffix.lower()
    if ext == ".docx":
        _run_docx(input_path, out_dir)
    elif ext == ".pptx":
        _run_pptx(input_path, out_dir)
    elif ext == ".xlsx":
        _run_xlsx(input_path, out_dir)
    else:
        raise ValueError(f"unsupported ext: {ext}")


def run_legacy(input_path: Path, out_dir: Path) -> None:
    ext = input_path.suffix.lower()
    target_ext_map = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx", ".xlsb": ".xlsx"}
    target_ext = target_ext_map[ext]

    converted_in_out = out_dir / f"_converted{target_ext}"
    _convert_legacy(input_path, converted_in_out)

    tool_extra = f"(레거시 {ext} → {target_ext} 변환 후 처리)"
    if target_ext == ".docx":
        _run_docx(converted_in_out, out_dir, source_name_override=input_path.name, tool_extra=tool_extra)
    elif target_ext == ".pptx":
        _run_pptx(converted_in_out, out_dir, source_name_override=input_path.name, tool_extra=tool_extra)
    elif target_ext == ".xlsx":
        _run_xlsx(converted_in_out, out_dir, source_name_override=input_path.name, tool_extra=tool_extra)


# ====================================================================
# DOCX
# ====================================================================

def _run_docx(
    input_path: Path,
    out_dir: Path,
    *,
    source_name_override: Optional[str] = None,
    tool_extra: str = "",
) -> None:
    pages_dir = out_dir / "pages"
    images_dir = out_dir / "images"

    # COM Word → 임시 PDF → PyMuPDF 로 페이지별 PNG + MD.
    with _common.com_lock(), _common.temp_workdir() as tmp:
        tmp_pdf = tmp / "out.pdf"
        _word_export_pdf(input_path, tmp_pdf)
        _common.mkdir(pages_dir)
        page_summaries = _render_pdf_pages(tmp_pdf, pages_dir)

    attachment_links = _extract_zip_media(
        input_path,
        out_dir,
        media_zip_prefix="word/media/",
        embed_zip_prefix="word/embeddings/",
        images_dir=images_dir,
    )

    source_name, source_size = _source_meta(input_path, out_dir, source_name_override)
    macro_modules = _extract_macros(_source_path(out_dir, source_name), out_dir)

    sections: dict[str, list[str]] = {}
    if page_summaries:
        sections[f"페이지 (총 {len(page_summaries)}개)"] = page_summaries
    if macro_modules:
        sections[f"VBA 매크로 (총 {len(macro_modules)}개)"] = [f"`macros/{m}`" for m in macro_modules]

    _common.write_readme(
        out_dir,
        source_name=source_name,
        source_size=source_size,
        tool=("COM Word + PyMuPDF + ZIP " + tool_extra).strip(),
        loss_notes="서식(폰트/색/볼드)·정확한 페이지 레이아웃은 PNG 안에서만 보존. 매크로(VBA)는 macros/ 로 별도 추출.",
        sections=sections or None,
        attachments=attachment_links,
    )


# ====================================================================
# PPTX
# ====================================================================

def _run_pptx(
    input_path: Path,
    out_dir: Path,
    *,
    source_name_override: Optional[str] = None,
    tool_extra: str = "",
) -> None:
    _common.ensure_pip("pptx", "python-pptx")
    from pptx import Presentation

    slides_dir = out_dir / "slides"
    charts_dir = out_dir / "charts"

    prs = Presentation(_common.long_str(input_path))
    slide_titles: list[tuple[str, str]] = []  # (idx, safe_title)
    slide_summaries: list[str] = []
    slide_has_notes: dict[str, bool] = {}
    slide_charts: dict[str, list[str]] = {}  # idx -> chart filenames

    _common.mkdir(slides_dir)
    for i, slide in enumerate(prs.slides, start=1):
        idx = f"{i:02d}"
        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
        if not title:
            title = f"슬라이드{i}"
        safe_title = _common.slugify_filename(title, max_len=40)
        slide_titles.append((idx, safe_title))

        # 슬라이드 텍스트 (python-pptx)
        text_lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run_.text for run_ in para.runs)
                if line.strip():
                    text_lines.append(line)
        _common.write_text(slides_dir / f"{idx}_{safe_title}.md", "\n".join(text_lines))

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text or ""
            if notes_text.strip():
                _common.write_text(
                    slides_dir / f"{idx}_{safe_title}.notes.md",
                    notes_text,
                )
                slide_has_notes[idx] = True

        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if shape.has_chart:
                # 차트 데이터 추출은 best-effort (비표준 차트면 못 뺄 수 있음).
                try:
                    data = _extract_pptx_chart_data(shape.chart)
                    _common.mkdir(charts_dir)
                    chart_filename = f"slide{i:02d}_chart{shape_idx:02d}.data.json"
                    _common.write_text(
                        charts_dir / chart_filename,
                        json.dumps(data, ensure_ascii=False, indent=2),
                    )
                    slide_charts.setdefault(idx, []).append(chart_filename)
                except Exception:
                    pass

        # 슬라이드별 산출물 풀목록
        parts = [f"`slides/{idx}_{safe_title}.png`", "`.md`"]
        if slide_has_notes.get(idx):
            parts.append("`.notes.md`")
        chart_refs = slide_charts.get(idx, [])
        if chart_refs:
            chart_str = ", ".join(f"`charts/{c}`" for c in chart_refs)
            parts.append(f"(차트: {chart_str})")
        slide_summaries.append(" ".join(parts))

    # COM PowerPoint 의 Slide.Export 로 슬라이드별 PNG 직접 출력. 임시 폴더에서 만든 후 long-path-safe copy.
    with _common.com_lock(), _common.temp_workdir() as tmp:
        _powerpoint_export_slides(input_path, tmp, slide_titles)
        for idx, safe_title in slide_titles:
            tmp_png = tmp / f"{idx}_{safe_title}.png"
            if tmp_png.exists():
                _common.copy(tmp_png, slides_dir / f"{idx}_{safe_title}.png")

    # pptx 의 시각은 슬라이드 PNG 에 모두 포함 → images/ 는 만들지 않음 (embeddings 만 추출).
    attachment_links = _extract_zip_media(
        input_path,
        out_dir,
        media_zip_prefix="ppt/media/",
        embed_zip_prefix="ppt/embeddings/",
    )

    source_name, source_size = _source_meta(input_path, out_dir, source_name_override)
    macro_modules = _extract_macros(_source_path(out_dir, source_name), out_dir)

    sections: dict[str, list[str]] = {}
    if slide_summaries:
        sections[f"슬라이드 (총 {len(slide_summaries)}개)"] = slide_summaries
    if macro_modules:
        sections[f"VBA 매크로 (총 {len(macro_modules)}개)"] = [f"`macros/{m}`" for m in macro_modules]

    _common.write_readme(
        out_dir,
        source_name=source_name,
        source_size=source_size,
        tool=("python-pptx + COM PowerPoint + ZIP " + tool_extra).strip(),
        loss_notes="애니메이션·슬라이드 전환·정확한 폰트는 미보존. 시각은 슬라이드별 PNG 로, 차트 데이터는 charts/*.data.json 으로 보존.",
        sections=sections or None,
        attachments=attachment_links,
    )


# ====================================================================
# XLSX
# ====================================================================

def _run_xlsx(
    input_path: Path,
    out_dir: Path,
    *,
    source_name_override: Optional[str] = None,
    tool_extra: str = "",
) -> None:
    _common.ensure_pip("openpyxl")
    from openpyxl import load_workbook
    from openpyxl.worksheet.worksheet import Worksheet

    sheets_dir = out_dir / "sheets"
    charts_dir = out_dir / "charts"
    sheet_summaries: list[str] = []
    sheet_names: list[tuple[str, str, str]] = []  # (idx, safe_name, raw_name)
    sheet_charts: dict[str, list[str]] = {}  # idx -> chart filenames
    sheet_formula_count: dict[str, int] = {}
    sheet_dims: dict[str, tuple[int, int]] = {}

    wb_values = load_workbook(_common.long_str(input_path), data_only=True)
    wb_formulas = load_workbook(_common.long_str(input_path), data_only=False)
    try:
        _common.mkdir(sheets_dir)
        # openpyxl 의 sheetnames 는 chart sheet 도 포함하지만 COM Worksheets() 는 일반 시트만.
        # 양쪽 인덱스 mismatch 피하려면 일반 Worksheet 만 처리하고, lookup 은 raw name 으로.
        idx_counter = 0
        for name in wb_values.sheetnames:
            if not isinstance(wb_values[name], Worksheet):
                continue
            idx_counter += 1
            idx = f"{idx_counter:02d}"
            safe_name = _common.slugify_filename(name, max_len=40)
            sheet_names.append((idx, safe_name, name))

        # COM Excel 호출: 데이터 영역 → ChartObject + Range.CopyPicture → 시트별 PNG.
        # 시트별 (last_row, last_col) 도 같이 반환되어 .md/.formulas.json 이 같은 데이터 영역으로 통일됨.
        with _common.com_lock():
            sheet_ranges = _excel_export_sheet_pngs(input_path, sheets_dir, sheet_names)

        for idx, safe_name, raw_name in sheet_names:
            ws_v = wb_values[raw_name]
            ws_f = wb_formulas[raw_name]

            # COM Find 결과가 있으면 그 범위, 없으면 openpyxl max_row/max_column fallback.
            last_row, last_col = sheet_ranges.get(raw_name, (ws_v.max_row, ws_v.max_column))
            sheet_dims[idx] = (last_row, last_col)

            md_lines = _sheet_to_md(ws_v, last_row, last_col)
            _common.write_text(sheets_dir / f"{idx}_{safe_name}.md", "\n".join(md_lines))

            formulas: dict[str, str] = {}
            if last_row >= 1 and last_col >= 1:
                for row in ws_f.iter_rows(min_row=1, max_row=last_row, min_col=1, max_col=last_col):
                    for cell in row:
                        v = cell.value
                        if isinstance(v, str) and v.startswith("="):
                            formulas[cell.coordinate] = v
            if formulas:
                _common.write_text(
                    sheets_dir / f"{idx}_{safe_name}.formulas.json",
                    json.dumps(formulas, ensure_ascii=False, indent=2),
                )
            sheet_formula_count[idx] = len(formulas)

            for chart_idx, chart in enumerate(getattr(ws_f, "_charts", []), start=1):
                # 차트 데이터 추출은 best-effort (비표준 차트면 못 뺄 수 있음).
                try:
                    data = _extract_openpyxl_chart_data(chart)
                    _common.mkdir(charts_dir)
                    chart_filename = f"sheet{idx}_chart{chart_idx:02d}.data.json"
                    _common.write_text(
                        charts_dir / chart_filename,
                        json.dumps(data, ensure_ascii=False, indent=2),
                    )
                    sheet_charts.setdefault(idx, []).append(chart_filename)
                except Exception:
                    pass
    finally:
        wb_values.close()
        wb_formulas.close()

    # 시트 PNG 는 데이터 영역(Find 범위) 만 캡처 → 데이터 영역 밖 이미지는 누락될 수 있음 →
    # raw 이미지를 시트+셀 위치 정보 포함해서 별도 보존.
    sheet_images = _extract_xlsx_images_with_position(input_path, out_dir, sheet_names)
    attachment_links = _extract_zip_media(
        input_path,
        out_dir,
        media_zip_prefix="xl/media/",
        embed_zip_prefix="xl/embeddings/",
    )

    # 시트별 산출물 풀목록 (모든 시트 처리 끝난 뒤 sheet_images 매핑까지 합쳐서)
    for idx, safe_name, raw_name in sheet_names:
        last_row, last_col = sheet_dims.get(idx, (0, 0))
        formula_n = sheet_formula_count.get(idx, 0)
        parts = [f"`sheets/{idx}_{safe_name}.png`", "`.md`"]
        if formula_n:
            parts.append("`.formulas.json`")
        chart_refs = sheet_charts.get(idx, [])
        if chart_refs:
            parts.append("(차트: " + ", ".join(f"`charts/{c}`" for c in chart_refs) + ")")
        img_refs = sheet_images.get(raw_name, [])
        if img_refs:
            parts.append("(이미지: " + ", ".join(f"`images/{n}`" for n in img_refs) + ")")
        meta = f"({last_row}행×{last_col}열"
        if formula_n:
            meta += f", 수식 {formula_n}개"
        meta += ")"
        sheet_summaries.append(" ".join(parts) + " " + meta)

    source_name, source_size = _source_meta(input_path, out_dir, source_name_override)
    macro_modules = _extract_macros(_source_path(out_dir, source_name), out_dir)

    sections: dict[str, list[str]] = {}
    if sheet_summaries:
        sections[f"시트 (총 {len(sheet_summaries)}개)"] = sheet_summaries
    if macro_modules:
        sections[f"VBA 매크로 (총 {len(macro_modules)}개)"] = [f"`macros/{m}`" for m in macro_modules]

    _common.write_readme(
        out_dir,
        source_name=source_name,
        source_size=source_size,
        tool=("openpyxl + COM Excel + ZIP " + tool_extra).strip(),
        loss_notes="셀 서식·조건부 서식·데이터 검증 규칙은 미보존. 시각은 시트별 PNG 로, 표 구조는 .md 로, 셀 수식은 .formulas.json 으로 보존.",
        sections=sections or None,
        attachments=attachment_links,
    )


# ====================================================================
# COM 헬퍼 — 모두 office_worker.py subprocess 로 격리.
# 이유: COM IUnknown::Release RPC release 가 Office 프로세스 cleanup 비동기성 때문에 60s 까지 wait 가능.
# 별도 process 로 격리하면 process 종료 = OS 가 COM cleanup 강제 → wait 없음 + hang 시 timeout kill.
# ====================================================================

_WORKER_PATH = Path(__file__).parent / "office_worker.py"


def _run_worker(cmd: str, *args: str, timeout: float, capture_stdout: bool = False) -> str:
    """office_worker.py subprocess 호출. timeout 초과 시 kill 후 throw.

    호출자는 com_lock 으로 sequential 보장. worker 자체는 단독 process 라 인스턴스 격리됨.
    """
    import subprocess
    proc = subprocess.run(
        [sys.executable, str(_WORKER_PATH), cmd, *args],
        capture_output=True, text=True, timeout=timeout, encoding="utf-8",
    )
    if proc.returncode != 0:
        raise RuntimeError(
            f"office_worker '{cmd}' failed (rc={proc.returncode}): {proc.stderr.strip()}"
        )
    return proc.stdout if capture_stdout else ""


def _word_export_pdf(input_path: Path, output_pdf: Path) -> None:
    _run_worker("word_pdf", str(input_path), str(output_pdf), timeout=600)


def _powerpoint_export_slides(
    input_path: Path,
    out_dir: Path,
    slide_titles: list[tuple[str, str]],
) -> None:
    """슬라이드별 PNG 직접 export. worker subprocess 로 격리."""
    _run_worker(
        "ppt_png", str(input_path), str(out_dir), json.dumps(slide_titles),
        timeout=600,
    )


def _excel_export_sheet_pngs(
    input_path: Path,
    sheets_dir: Path,
    sheet_names: list[tuple[str, str, str]],
) -> dict[str, tuple[int, int]]:
    """시트별 PNG 생성 + (last_row, last_col) 매핑 반환.

    호출자에서 sheetProtection strip 사본 만들고 worker 에 그 사본 path 만 넘김.
    Excel COM 자체 작업은 worker subprocess.
    """
    with _common.temp_workdir() as tmp:
        unprotected = tmp / "_unprotected.xlsx"
        _xlsx_strip_protection(input_path, unprotected)
        result = _run_worker(
            "excel_sheets", str(unprotected), str(sheets_dir), json.dumps(sheet_names),
            timeout=600, capture_stdout=True,
        )
    raw = json.loads(result) if result.strip() else {}
    return {k: tuple(v) for k, v in raw.items()}


def _xlsx_strip_protection(src: Path, dst: Path) -> None:
    """xlsx ZIP 안 sheetProtection / workbookProtection xml 노드 제거한 사본 생성.

    Excel COM 의 ChartObjects.Add 가 보호 시트에서 차단되는 이슈 회피용.
    비번 hash 우회 불필요 — xlsx 의 보호는 xml 노드가 본체. 다른 part 는 그대로 → 차트/매크로/이미지 보존.
    """
    import re
    # attribute value 에 '/' 들어갈 수 있음 (hashValue base64) → '/' 제외하면 안 됨.
    # '>' 만 제외하고 lazy 로 첫 '/>' 까지.
    sheet_re = re.compile(rb'<sheetProtection\b[^>]*?/>')
    wb_re = re.compile(rb'<workbookProtection\b[^>]*?/>')
    with zipfile.ZipFile(_common.long_str(src), "r") as zin, \
         zipfile.ZipFile(_common.long_str(dst), "w", zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info)
            name = info.filename
            if name.startswith("xl/worksheets/sheet") and name.endswith(".xml"):
                data = sheet_re.sub(b"", data)
            elif name == "xl/workbook.xml":
                data = wb_re.sub(b"", data)
            zout.writestr(info, data)


def _convert_legacy(input_path: Path, target_path: Path) -> None:
    """input_path (.doc/.ppt/.xls/.xlsb) 를 신형 형식으로 변환해 target_path 에 저장.

    임시 폴더에 변환 후 long-path-safe copy. Office COM 작업은 worker subprocess.
    """
    target_ext = target_path.suffix
    with _common.com_lock(), _common.temp_workdir() as tmp:
        tmp_target = tmp / (input_path.stem + target_ext)
        _run_worker(
            "convert_legacy", input_path.suffix.lower(), str(input_path), str(tmp_target),
            timeout=600,
        )
        _common.copy(tmp_target, target_path)


# ====================================================================
# 공용
# ====================================================================

def _render_pdf_pages(pdf_path: Path, pages_dir: Path) -> list[str]:
    """PDF → 페이지별 PNG + MD."""
    _common.ensure_pip("fitz", "PyMuPDF")
    import fitz

    summaries: list[str] = []
    doc = fitz.open(_common.long_str(pdf_path))
    try:
        for i, page in enumerate(doc, start=1):
            idx = f"{i:03d}"
            text = page.get_text("text") or ""
            _common.write_text(pages_dir / f"{idx}.md", text)
            pix = page.get_pixmap(dpi=300)
            pix.save(_common.long_str(pages_dir / f"{idx}.png"))
            summaries.append(f"`pages/{idx}.png` (시각) — `.md` ({len(text)}자)")
    finally:
        doc.close()
    return summaries


def _source_meta(
    input_path: Path,
    out_dir: Path,
    source_name_override: Optional[str],
) -> tuple[str, int]:
    if source_name_override:
        src_ext = Path(source_name_override).suffix.lstrip(".")
        src_path = out_dir / f"_source.{src_ext}"
        return source_name_override, src_path.stat().st_size
    return input_path.name, input_path.stat().st_size


def _source_path(out_dir: Path, source_name: str) -> Path:
    """out_dir 안 보존된 원본 _source.<ext> 경로."""
    ext = Path(source_name).suffix.lstrip(".")
    return out_dir / f"_source.{ext}"


def _extract_macros(input_path: Path, out_dir: Path) -> list[str]:
    """OLE/OOXML 파일에서 VBA 매크로 추출. macros/<모듈명>.vba 로 저장.

    추출된 모듈 파일명 list 반환 (예: ["Module1.vba", "ThisWorkbook.vba"]).
    매크로 없으면 빈 list.
    """
    _common.ensure_pip("oletools")
    from oletools.olevba import VBA_Parser

    module_files: list[str] = []
    parser = VBA_Parser(_common.long_str(input_path))
    try:
        if not parser.detect_vba_macros():
            return module_files
        macros_dir = out_dir / "macros"
        _common.mkdir(macros_dir)
        for (_filename, stream_path, vba_filename, vba_code) in parser.extract_macros():
            module_name = vba_filename or stream_path or "module"
            stem = Path(module_name).stem or "module"
            dst = _common.unique_path(macros_dir, f"{stem}.vba")
            _common.write_text(dst, vba_code or "")
            module_files.append(dst.name)
        return module_files
    finally:
        parser.close()


def _extract_xlsx_images_with_position(
    input_path: Path,
    out_dir: Path,
    sheet_names: list[tuple[str, str, str]],
) -> dict[str, list[str]]:
    """xlsx 의 이미지를 시트+셀 위치 정보 포함 파일명으로 images/ 에 저장.

    파일명 패턴: `<sheet_safe_name>_<cell_addr>.<ext>` (예: `BOA_E5.png`).
    같은 셀에 여러 이미지면 unique_path 가 _1, _2 suffix.
    이미지 데이터/anchor 추출 실패 시 그 이미지만 skip (전체 fail 안 함).
    """
    _common.ensure_pip("openpyxl")
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    images_dir = out_dir / "images"
    name_to_safe = {raw: safe for _idx, safe, raw in sheet_names}
    sheet_images: dict[str, list[str]] = {}  # raw_name -> [filename, ...]

    wb = load_workbook(_common.long_str(input_path))
    try:
        for sheet_name in wb.sheetnames:
            if sheet_name not in name_to_safe:
                continue  # chart sheet 등 일반 worksheet 아닌 경우 skip
            ws = wb[sheet_name]
            safe_name = name_to_safe[sheet_name]
            for img in getattr(ws, "_images", []):
                # anchor 위치 추출 (OneCellAnchor / TwoCellAnchor 둘 다 _from 갖음)
                cell_addr = "anchor_unknown"
                try:
                    anchor = img.anchor
                    marker = getattr(anchor, "_from", None) or getattr(anchor, "from", None)
                    if marker is not None:
                        col = marker.col + 1  # 0-based → 1-based
                        row = marker.row + 1
                        cell_addr = f"{get_column_letter(col)}{row}"
                except Exception:
                    pass

                # 이미지 raw 데이터 추출 (openpyxl 내부 API).
                try:
                    data = img._data()
                except Exception:
                    continue
                if not data:
                    continue

                # 확장자: magic bytes 로 추정.
                ext = "bin"
                if data.startswith(b"\x89PNG"):
                    ext = "png"
                elif data.startswith(b"\xff\xd8"):
                    ext = "jpg"
                elif data.startswith(b"GIF8"):
                    ext = "gif"
                elif data.startswith(b"BM"):
                    ext = "bmp"
                elif data[:4] == b"RIFF" and data[8:12] == b"WEBP":
                    ext = "webp"

                _common.mkdir(images_dir)
                base = f"{safe_name}_{cell_addr}.{ext}"
                dst = _common.unique_path(images_dir, base)
                _common.write_bytes(dst, data)
                sheet_images.setdefault(sheet_name, []).append(dst.name)
    finally:
        wb.close()
    return sheet_images


def _extract_zip_media(
    input_path: Path,
    out_dir: Path,
    *,
    media_zip_prefix: str,
    embed_zip_prefix: str,
    images_dir: Optional[Path] = None,
) -> list[str]:
    """OOXML(docx/pptx/xlsx) ZIP 안의 media/ 와 embeddings/ 에서 이미지/첨부 추출.

    images_dir=None 이면 media/ 추출 skip (PDF 가 시각 책임지는 형식: docx, pptx).
    xlsx 는 별도 _extract_xlsx_images_with_position 으로 위치 정보 포함 추출.
    """
    attachment_links: list[str] = []
    attachments_dir = out_dir / "attachments"

    with zipfile.ZipFile(_common.long_str(input_path), "r") as zf:
        for info in zf.infolist():
            name = info.filename
            if info.is_dir():
                continue
            if images_dir is not None and name.startswith(media_zip_prefix):
                _common.mkdir(images_dir)
                base = Path(name).name
                dst = _common.unique_path(images_dir, base)
                with zf.open(info) as f:
                    _common.write_bytes(dst, f.read())
            elif name.startswith(embed_zip_prefix):
                _common.mkdir(attachments_dir)
                base = Path(name).name
                dst = _common.unique_path(attachments_dir, base)
                with zf.open(info) as f:
                    _common.write_bytes(dst, f.read())
                recursed = maybe_recurse_attachment(dst, attachments_dir)
                if recursed is not None:
                    os.unlink(_common.long_str(dst))
                    attachment_links.append(f"attachments/{recursed.name}/")
                else:
                    attachment_links.append(f"attachments/{dst.name}")
    return attachment_links


def _sheet_to_md(ws, last_row: int, last_col: int) -> list[str]:
    """openpyxl Worksheet 의 (1,1)~(last_row,last_col) 범위를 마크다운 표 라인으로."""
    if last_row < 1 or last_col < 1:
        return ["(빈 시트)"]

    rows: list[list[str]] = []
    for row in ws.iter_rows(
        min_row=1, max_row=last_row, min_col=1, max_col=last_col, values_only=True
    ):
        rows.append(["" if v is None else str(v) for v in row])
    if not rows or not any(any(c for c in r) for r in rows):
        return ["(빈 시트)"]

    header = rows[0]
    md_lines: list[str] = []
    md_lines.append("| " + " | ".join(_md_escape(c) for c in header) + " |")
    md_lines.append("| " + " | ".join("---" for _ in header) + " |")
    for row in rows[1:]:
        padded = list(row) + [""] * (len(header) - len(row))
        md_lines.append("| " + " | ".join(_md_escape(c) for c in padded[: len(header)]) + " |")
    return md_lines


def _md_escape(s: str) -> str:
    return s.replace("|", "\\|").replace("\n", " ")


def _extract_pptx_chart_data(chart) -> dict:
    data: dict = {"type": str(getattr(chart, "chart_type", "")), "categories": [], "series": []}
    if chart.plots:
        cats = list(chart.plots[0].categories)
        data["categories"] = [str(c) for c in cats]
    for plot in chart.plots:
        for series in plot.series:
            values = list(series.values) if hasattr(series, "values") else []
            data["series"].append({
                "name": str(getattr(series, "name", "")),
                "values": [v for v in values],
            })
    return data


def _extract_openpyxl_chart_data(chart) -> dict:
    data: dict = {"type": type(chart).__name__, "title": "", "series": []}
    if chart.title:
        data["title"] = str(chart.title)
    for s in chart.series:
        ref = ""
        if s.val and s.val.numRef:
            ref = str(s.val.numRef.f)
        data["series"].append({"ref": ref})
    return data
