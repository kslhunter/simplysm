#!/usr/bin/env python3
"""Extract text and images from PPTX files with per-slide coordinates."""

import sys
from _common import (
    setup_encoding, make_output_paths, print_header, save_image,
    ext_from_content_type, print_image_summary, run_cli,
)

setup_encoding()


def emu_to_inches(emu):
    """Convert EMU to inches (1 decimal place)."""
    if emu is None:
        return "?"
    return f"{emu / 914400:.1f}"


def _pos(shape):
    return f"(left={emu_to_inches(shape.left)}\", top={emu_to_inches(shape.top)}\")"


def extract(file_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    fp, out_dir = make_output_paths(file_path)
    img_idx = 0

    print_header(fp)

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"## Slide {slide_num}\n")

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img_idx += 1
                ext = ext_from_content_type(shape.image.content_type)
                img_path = save_image(out_dir, img_idx, shape.image.blob, ext)
                print(f"[IMG] {_pos(shape)} {img_path}")

            elif hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip().replace("\n", "\n       ")
                print(f"[TXT] {_pos(shape)} {text}")

        print()

    print_image_summary(img_idx, out_dir)


if __name__ == "__main__":
    run_cli(extract, "extract_pptx.py", {"python-pptx": "pptx"})
