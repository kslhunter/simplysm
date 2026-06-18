"""Office (docx/pptx/xlsx) + 레거시 (doc/ppt/xls/xlsb) COM 핸들러.

시각 산출물은 PNG, 텍스트/구조 산출물은 형식별로:
- docx → pages/<NNN>.png + content.jsonl + pages.meta.json
- pptx → slides/<NN>_<title>.png + .jsonl (+ .notes.md 있을 때) (슬라이드별)
- xlsx → sheets/<NN>_<name>.png + .jsonl (시트별) + workbook.meta.json

xlsx jsonl 한 줄 = 한 행. 좌표는 행번호(`r`)·열문자 키로 명시. 값·수식·시트 메타 통합.

Office COM 호출은 office_worker.py subprocess 로 격리 (cleanup race 회피).
이 모듈 (office_com.py) 은 호출자 + Office 외 작업 (jsonl 직렬화, ZIP strip, 매크로 추출, README 생성).
원칙: 처리 실패는 묻지 않고 그대로 throw. try/finally 는 락/임시 폴더 cleanup 에만 사용.
"""
from __future__ import annotations

import json
import os
import re
import sys
import zipfile
from datetime import date, datetime, time
from pathlib import Path
from typing import Any, Optional

from . import _common
from .dispatch import maybe_recurse_attachment


# ====================================================================
# 진입점
# ====================================================================

def run(input_path: Path, out_dir: Path) -> None:
    ext = input_path.suffix.lower()
    if ext == ".docx":
        _run_docx(input_path, out_dir)
    elif ext == ".pptx":
        _run_pptx(input_path, out_dir)
    elif ext == ".xlsx":
        _run_xlsx(input_path, out_dir)
    else:
        raise ValueError(f"unsupported ext: {ext}")


def run_legacy(input_path: Path, out_dir: Path) -> None:
    """레거시 (.doc/.ppt/.xls/.xlsb) → 신형 변환 후 처리.

    `_converted.<ext>` 는 임시 폴더에서만 처리하고 산출 폴더(out_dir)에는 잔존시키지 않음.
    """
    ext = input_path.suffix.lower()
    target_ext_map = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx", ".xlsb": ".xlsx"}
    target_ext = target_ext_map[ext]

    tool_extra = f"(레거시 {ext} → {target_ext} 변환 후 처리)"
    with _common.temp_workdir() as tmp:
        converted_path = tmp / f"_converted{target_ext}"
        _convert_legacy(input_path, converted_path)

        if target_ext == ".docx":
            _run_docx(converted_path, out_dir, source_name_override=input_path.name, tool_extra=tool_extra)
        elif target_ext == ".pptx":
            _run_pptx(converted_path, out_dir, source_name_override=input_path.name, tool_extra=tool_extra)
        elif target_ext == ".xlsx":
            _run_xlsx(converted_path, out_dir, source_name_override=input_path.name, tool_extra=tool_extra)


# ====================================================================
# DOCX
# ====================================================================

def _run_docx(
    input_path: Path,
    out_dir: Path,
    *,
    source_name_override: Optional[str] = None,
    tool_extra: str = "",
) -> None:
    """python-docx 로 구조 추출 → content.jsonl 단일 시퀀스. 페이지 단위 폐기.

    PNG 는 fitz PDF 경유로 시각 검증용 유지. pages.meta.json 으로 페이지↔노드 best-effort 매핑.
    """
    _common.ensure_pip("docx", "python-docx")

    pages_dir = out_dir / "pages"
    images_dir = out_dir / "images"

    # 1. python-docx 구조 추출
    nodes, counts = _docx_extract_nodes(input_path)

    # content.jsonl
    lines: list[str] = [json.dumps({"_meta": counts}, ensure_ascii=False)]
    for n in nodes:
        lines.append(json.dumps(n, ensure_ascii=False, default=_json_default))
    _common.write_text(out_dir / "content.jsonl", "\n".join(lines))

    # 2. fitz PDF 경유 PNG + pages.meta.json (페이지↔노드 매핑 best-effort)
    with _common.com_lock(), _common.temp_workdir() as tmp:
        tmp_pdf = tmp / "out.pdf"
        _word_export_pdf(input_path, tmp_pdf)
        _common.mkdir(pages_dir)
        page_count = _docx_pages_from_pdf(tmp_pdf, pages_dir, out_dir, nodes)

    attachment_links = _extract_zip_media(
        input_path,
        out_dir,
        media_zip_prefix="word/media/",
        embed_zip_prefix="word/embeddings/",
        images_dir=images_dir,
    )

    source_name, source_size = _source_meta(input_path, out_dir, source_name_override)
    macro_modules = _extract_macros(_source_path(out_dir, source_name), out_dir)

    sections: dict[str, list[str]] = {}
    summary = (
        f"노드 {counts['nodes']}개 "
        f"(heading {counts['headings']}·para {counts['paragraphs_plain']}·"
        f"bullet {counts['bullets']}·table_cell {counts['table_cells']}·image {counts['images']})"
    )
    content_items = [f"`content.jsonl` — {summary}"]
    if page_count:
        content_items.append(f"`pages.meta.json` — PNG ↔ 노드 매핑 ({page_count}페이지)")
    sections["콘텐츠"] = content_items
    if macro_modules:
        sections[f"VBA 매크로 (총 {len(macro_modules)}개)"] = [f"`macros/{m}`" for m in macro_modules]

    _common.write_readme(
        out_dir,
        source_name=source_name,
        source_size=source_size,
        tool=("python-docx + COM Word + PyMuPDF + ZIP " + tool_extra).strip(),
        loss_notes=(
            "서식(폰트/색/볼드)·정확한 페이지 레이아웃은 PNG 안에서만 보존. "
            "구조는 content.jsonl 단일 시퀀스(heading/para/bullet/table_cell/image), "
            "PNG↔노드 매핑은 pages.meta.json. 매크로(VBA)는 macros/ 로 별도 추출."
        ),
        sections=sections or None,
        attachments=attachment_links,
    )


def _docx_extract_nodes(input_path: Path) -> tuple[list[dict], dict[str, int]]:
    """python-docx 로 body 시퀀스(paragraph/table) 순회 → jsonl 노드 리스트."""
    from docx import Document
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    IMG_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
    doc = Document(_common.long_str(input_path))

    # image relationship: rid → 'images/<basename>'
    img_rels: dict[str, str] = {}
    for rid, rel in doc.part.rels.items():
        if rel.reltype == IMG_RELTYPE:
            try:
                basename = Path(rel.target_ref).name
                img_rels[rid] = f"images/{basename}"
            except Exception:
                continue

    nodes: list[dict] = []
    counts = {
        "nodes": 0,
        "headings": 0,
        "paragraphs_plain": 0,
        "bullets": 0,
        "tables": 0,
        "table_cells": 0,
        "images": 0,
    }
    table_idx = 0
    node_idx = 0

    for elem in doc.element.body.iterchildren():
        tag = elem.tag
        if tag == qn("w:p"):
            para = Paragraph(elem, doc)
            text = para.text or ""
            style_name = para.style.name if para.style else ""
            heading_level = _docx_heading_level(style_name)
            bullet_level = _docx_bullet_level(para)
            image_rids = _docx_inline_image_rids(para)
            hyperlinks = _docx_paragraph_hyperlinks(para, doc)

            node: dict
            if heading_level is not None:
                node = {"node": node_idx, "type": "heading", "level": heading_level, "text": text}
                counts["headings"] += 1
            elif bullet_level is not None:
                node = {"node": node_idx, "type": "bullet", "level": bullet_level, "text": text}
                counts["bullets"] += 1
            else:
                # 빈 paragraph 도 원본 정보 → 노드로 보존 (text="")
                node = {"node": node_idx, "type": "para", "text": text}
                counts["paragraphs_plain"] += 1

            if hyperlinks:
                node["hyperlinks"] = hyperlinks

            nodes.append(node)
            node_idx += 1

            for rid in image_rids:
                ref = img_rels.get(rid)
                if ref:
                    nodes.append({"node": node_idx, "type": "image", "ref": ref})
                    counts["images"] += 1
                    node_idx += 1

        elif tag == qn("w:tbl"):
            table_obj = Table(elem, doc)
            table_idx += 1
            counts["tables"] += 1
            seen_tc: set[int] = set()
            for r, row in enumerate(table_obj.rows, start=1):
                for c, cell in enumerate(row.cells, start=1):
                    tc_id = id(cell._tc)
                    if tc_id in seen_tc:
                        # gridSpan 으로 같은 row 안 colspan 중복 노출 — origin 의 colspan 에 표기됨
                        continue
                    seen_tc.add(tc_id)
                    vm = _docx_cell_vmerge(cell)
                    if vm == "continue":
                        # vMerge continue cell — origin 의 rowspan 영역. skip.
                        continue
                    cell_text = cell.text or ""  # 원본 그대로 (strip X)
                    colspan = _docx_cell_colspan(cell)
                    cell_node = {
                        "node": node_idx,
                        "type": "table_cell",
                        "table_idx": table_idx,
                        "row": r,
                        "col": c,
                        "text": cell_text,
                    }
                    if colspan > 1:
                        cell_node["colspan"] = colspan
                    nodes.append(cell_node)
                    counts["table_cells"] += 1
                    node_idx += 1

    counts["nodes"] = node_idx
    return nodes, counts


def _docx_heading_level(style_name: str) -> Optional[int]:
    """python-docx 스타일명 → heading level. heading 아니면 None."""
    if not style_name:
        return None
    if style_name.startswith("Heading "):
        try:
            return int(style_name.split(" ")[1])
        except (ValueError, IndexError):
            return None
    if style_name == "Title":
        return 0
    return None


def _docx_bullet_level(para) -> Optional[int]:
    """paragraph 의 numbering ilvl 추출. bullet/numbered 아니면 None."""
    from docx.oxml.ns import qn

    pPr = para._element.find(qn("w:pPr"))
    if pPr is None:
        return None
    numPr = pPr.find(qn("w:numPr"))
    if numPr is None:
        return None
    ilvl_elem = numPr.find(qn("w:ilvl"))
    if ilvl_elem is None:
        return 0
    try:
        return int(ilvl_elem.get(qn("w:val")) or 0)
    except (ValueError, TypeError):
        return 0


_DRAWING_EMBED_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
_DRAWING_BLIP_TAG = "{http://schemas.openxmlformats.org/drawingml/2006/main}blip"


def _docx_inline_image_rids(para) -> list[str]:
    """paragraph 안 inline image relationship IDs."""
    from docx.oxml.ns import qn

    rids: list[str] = []
    for drawing in para._element.iter(qn("w:drawing")):
        for blip in drawing.iter(_DRAWING_BLIP_TAG):
            rid = blip.get(_DRAWING_EMBED_NS)
            if rid:
                rids.append(rid)
    return rids


_DOCX_R_ID_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
_DOCX_HYPERLINK_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"


def _docx_paragraph_hyperlinks(para, doc) -> list[dict]:
    """paragraph 안 hyperlink list: [{"text":"...", "url":"..."}, ...]"""
    from docx.oxml.ns import qn

    rels = doc.part.rels
    result: list[dict] = []
    for hl_elem in para._element.iter(qn("w:hyperlink")):
        rid = hl_elem.get(_DOCX_R_ID_NS)
        url = ""
        if rid and rid in rels:
            rel = rels[rid]
            if rel.reltype == _DOCX_HYPERLINK_RELTYPE:
                url = rel.target_ref
        # hyperlink 안 모든 w:t 텍스트 join
        hl_text = "".join((t.text or "") for t in hl_elem.iter(qn("w:t")))
        if hl_text or url:
            result.append({"text": hl_text, "url": url})
    return result


def _docx_cell_colspan(cell) -> int:
    """docx 표 셀의 colspan (gridSpan val). 기본 1."""
    from docx.oxml.ns import qn

    tcPr = cell._tc.find(qn("w:tcPr"))
    if tcPr is None:
        return 1
    gridSpan = tcPr.find(qn("w:gridSpan"))
    if gridSpan is None:
        return 1
    val = gridSpan.get(qn("w:val"))
    try:
        return int(val) if val else 1
    except (ValueError, TypeError):
        return 1


def _docx_cell_vmerge(cell) -> Optional[str]:
    """docx 표 셀의 vMerge 상태. 'restart' | 'continue' | None."""
    from docx.oxml.ns import qn

    tcPr = cell._tc.find(qn("w:tcPr"))
    if tcPr is None:
        return None
    vMerge = tcPr.find(qn("w:vMerge"))
    if vMerge is None:
        return None
    val = vMerge.get(qn("w:val"))
    return val if val else "continue"  # vMerge 요소 있고 val 없으면 continue


def _docx_pages_from_pdf(
    pdf_path: Path,
    pages_dir: Path,
    out_dir: Path,
    nodes: list[dict],
) -> int:
    """fitz PDF 경유 페이지별 PNG + pages.meta.json (페이지별 raw text 보존).

    nodes 와의 매핑은 fitz·python-docx 간 텍스트 분할 차이로 자동 추정 시 오매핑 위험 →
    raw text 만 보존. Claude 가 분석 시 페이지 text 와 content.jsonl 노드 text 를 직접 비교.
    """
    _common.ensure_pip("fitz", "PyMuPDF")
    import fitz

    pages_meta: dict[str, dict] = {}
    fdoc = fitz.open(_common.long_str(pdf_path))
    try:
        for i, page in enumerate(fdoc, start=1):
            idx = f"{i:03d}"
            pix = page.get_pixmap(dpi=300)
            pix.save(_common.long_str(pages_dir / f"{idx}.png"))
            text = page.get_text("text") or ""
            pages_meta[idx] = {"text": text}
    finally:
        fdoc.close()

    if pages_meta:
        _common.write_text(
            out_dir / "pages.meta.json",
            json.dumps(pages_meta, ensure_ascii=False, indent=2),
        )
    return len(pages_meta)


# ====================================================================
# PPTX
# ====================================================================

def _run_pptx(
    input_path: Path,
    out_dir: Path,
    *,
    source_name_override: Optional[str] = None,
    tool_extra: str = "",
) -> None:
    """python-pptx 로 구조 추출 → 슬라이드별 jsonl. 시각 순서 정렬 + pos EMU 좌표.

    노드 type: title·heading·para·bullet·table_cell·image·chart·shape.
    PNG 은 COM PowerPoint 의 Slide.Export 로 슬라이드별 직접 출력.
    """
    _common.ensure_pip("pptx", "python-pptx")
    from pptx import Presentation

    slides_dir = out_dir / "slides"
    charts_dir = out_dir / "charts"
    images_dir = out_dir / "images"

    prs = Presentation(_common.long_str(input_path))
    slide_w = int(prs.slide_width or 0)
    slide_h = int(prs.slide_height or 0)

    slide_titles: list[tuple[str, str]] = []  # (idx, safe_title)
    slide_summaries: list[str] = []
    slide_has_notes: dict[str, bool] = {}
    slide_charts: dict[str, list[str]] = {}  # idx -> chart filenames

    _common.mkdir(slides_dir)
    for i, slide in enumerate(prs.slides, start=1):
        idx = f"{i:02d}"
        title = _pptx_slide_title(slide)
        # title 없으면 idx 만 (자체 한국어 라벨 부착 X)
        safe_title = _common.slugify_filename(title, max_len=40) if title else ""
        slide_titles.append((idx, safe_title))

        nodes, chart_refs = _pptx_extract_slide_nodes(
            slide, i, charts_dir, images_dir,
        )
        # 원본 XML 순서 (shape_idx 순) 그대로 보존. 시각 순서는 pos 가 보존되어 있어
        # Claude 가 필요시 직접 정렬 가능.

        meta = {
            "_meta": {
                "slide": i,
                "title": title,
                "size": [slide_w, slide_h],
                "shapes": len(nodes),
            }
        }
        lines = [json.dumps(meta, ensure_ascii=False, default=_json_default)]
        for n in nodes:
            lines.append(json.dumps(n, ensure_ascii=False, default=_json_default))
        stem = _pptx_slide_stem(idx, safe_title)
        _common.write_text(slides_dir / f"{stem}.jsonl", "\n".join(lines))

        if chart_refs:
            slide_charts[idx] = chart_refs

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text or ""
            if notes_text:
                _common.write_text(
                    slides_dir / f"{stem}.notes.md",
                    notes_text,
                )
                slide_has_notes[idx] = True

        parts = [f"`slides/{stem}.png`", "`.jsonl`"]
        if slide_has_notes.get(idx):
            parts.append("`.notes.md`")
        if chart_refs:
            chart_str = ", ".join(f"`charts/{c}`" for c in chart_refs)
            parts.append(f"(차트: {chart_str})")
        slide_summaries.append(" ".join(parts))

    # COM PowerPoint 의 Slide.Export 로 슬라이드별 PNG 직접 출력. 임시 폴더에서 만든 후 long-path-safe copy.
    with _common.com_lock(), _common.temp_workdir() as tmp:
        _powerpoint_export_slides(input_path, tmp, slide_titles)
        for idx, safe_title in slide_titles:
            stem = _pptx_slide_stem(idx, safe_title)
            tmp_png = tmp / f"{stem}.png"
            if tmp_png.exists():
                _common.copy(tmp_png, slides_dir / f"{stem}.png")

    # pptx 의 시각은 슬라이드 PNG 에 모두 포함 → ZIP media 전체 복제 skip
    # (개별 picture shape 은 _pptx_extract_slide_nodes 에서 image ref 와 함께 저장됨).
    attachment_links = _extract_zip_media(
        input_path,
        out_dir,
        media_zip_prefix="ppt/media/",
        embed_zip_prefix="ppt/embeddings/",
    )

    source_name, source_size = _source_meta(input_path, out_dir, source_name_override)
    macro_modules = _extract_macros(_source_path(out_dir, source_name), out_dir)

    sections: dict[str, list[str]] = {}
    if slide_summaries:
        sections[f"슬라이드 (총 {len(slide_summaries)}개)"] = slide_summaries
    if macro_modules:
        sections[f"VBA 매크로 (총 {len(macro_modules)}개)"] = [f"`macros/{m}`" for m in macro_modules]

    _common.write_readme(
        out_dir,
        source_name=source_name,
        source_size=source_size,
        tool=("python-pptx + COM PowerPoint + ZIP " + tool_extra).strip(),
        loss_notes=(
            "애니메이션·슬라이드 전환·정확한 폰트는 미보존. "
            "시각은 슬라이드별 PNG, 구조는 슬라이드별 .jsonl(시각 순서·pos EMU 좌표), "
            "차트 데이터는 charts/*.data.json, picture shape 의 image 는 images/."
        ),
        sections=sections or None,
        attachments=attachment_links,
    )


def _pptx_slide_stem(idx: str, safe_title: str) -> str:
    """슬라이드 파일 stem. safe_title 빈 문자열이면 idx 만 (자체 라벨 부착 X)."""
    return f"{idx}_{safe_title}" if safe_title else idx


def _pptx_slide_title(slide) -> str:
    """슬라이드 title placeholder 텍스트. 없으면 빈 문자열. 원본 그대로 (strip X)."""
    try:
        title_shape = slide.shapes.title
        if title_shape is not None and title_shape.text:
            return title_shape.text
    except (AttributeError, ValueError):
        pass
    return ""


def _pptx_extract_slide_nodes(
    slide,
    slide_num: int,
    charts_dir: Path,
    images_dir: Path,
) -> tuple[list[dict], list[str]]:
    """슬라이드 안 shape → 노드 list + chart 파일 list.

    text_frame 의 paragraph 별로 노드 분리 (heading·para·bullet).
    표·차트·이미지는 각각 별도 노드.
    그 외 (autoshape·SmartArt·group) 은 shape 노드.
    """
    nodes: list[dict] = []
    chart_refs: list[str] = []

    title_shape = None
    try:
        title_shape = slide.shapes.title
    except (AttributeError, ValueError):
        title_shape = None

    for shape_idx, shape in enumerate(slide.shapes):
        pos = _pptx_shape_pos(shape)
        common = {
            "slide": slide_num,
            "pos": pos,
            "shape_idx": shape_idx,
        }

        # 표
        if getattr(shape, "has_table", False):
            try:
                table = shape.table
            except Exception:
                table = None
            if table is not None:
                table_idx = shape_idx + 1
                for r_idx, row in enumerate(table.rows, start=1):
                    for c_idx, cell in enumerate(row.cells, start=1):
                        cell_text = cell.text or ""  # 원본 그대로 (strip X)
                        nodes.append({
                            **common,
                            "type": "table_cell",
                            "table_idx": table_idx,
                            "row": r_idx,
                            "col": c_idx,
                            "text": cell_text,
                        })
                continue

        # 차트
        if getattr(shape, "has_chart", False):
            try:
                data = _extract_pptx_chart_data(shape.chart)
            except Exception:
                data = None
            chart_filename = f"slide{slide_num:02d}_chart{shape_idx + 1:02d}.data.json"
            if data is not None:
                _common.mkdir(charts_dir)
                _common.write_text(
                    charts_dir / chart_filename,
                    json.dumps(data, ensure_ascii=False, indent=2),
                )
                chart_refs.append(chart_filename)
            nodes.append({
                **common,
                "type": "chart",
                "ref": f"charts/{chart_filename}",
            })
            continue

        # 그림 (picture)
        if _pptx_is_picture(shape):
            ref = _pptx_save_picture(shape, slide_num, shape_idx, images_dir)
            node = {**common, "type": "image"}
            if ref:
                node["ref"] = ref
            nodes.append(node)
            continue

        # text_frame 보유 shape (placeholder·text box·autoshape with text)
        if getattr(shape, "has_text_frame", False):
            is_title = (title_shape is not None and shape == title_shape)
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                text = "".join(run.text for run in para.runs)
                hyperlinks = _pptx_run_hyperlinks(para)
                bullet_lvl = getattr(para, "level", 0) or 0

                base_node: dict
                if is_title and p_idx == 0:
                    base_node = {**common, "type": "title", "para_idx": p_idx, "text": text}
                elif bullet_lvl > 0:
                    base_node = {**common, "type": "bullet", "para_idx": p_idx,
                                 "level": bullet_lvl, "text": text}
                else:
                    base_node = {**common, "type": "para", "para_idx": p_idx, "text": text}
                if hyperlinks:
                    base_node["hyperlinks"] = hyperlinks
                nodes.append(base_node)
            continue

        # 그 외 (group·SmartArt·connector·autoshape 등)
        subtype = ""
        try:
            subtype = str(shape.shape_type)
        except Exception:
            pass
        nodes.append({
            **common,
            "type": "shape",
            "subtype": subtype,
        })

    return nodes, chart_refs


def _pptx_shape_pos(shape) -> list[int]:
    """shape 의 [left, top, width, height] EMU. 누락 시 0."""
    try:
        return [
            int(shape.left or 0),
            int(shape.top or 0),
            int(shape.width or 0),
            int(shape.height or 0),
        ]
    except (AttributeError, TypeError, ValueError):
        return [0, 0, 0, 0]


def _pptx_is_picture(shape) -> bool:
    """python-pptx shape 이 picture 인지. shape_type 또는 image 속성으로 판별."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return True
    except Exception:
        pass
    # placeholder picture 인 경우 shape_type 이 PLACEHOLDER 라 image 속성으로 보완
    try:
        _ = shape.image
        return True
    except Exception:
        return False


def _pptx_run_hyperlinks(para) -> list[dict]:
    """pptx paragraph 안 run 별 hyperlink list. 텍스트·URL."""
    result: list[dict] = []
    for run in para.runs:
        try:
            hl = run.hyperlink
            url = getattr(hl, "address", None)
        except Exception:
            url = None
        if url:
            result.append({"text": run.text or "", "url": url})
    return result


def _pptx_save_picture(
    shape, slide_num: int, shape_idx: int, images_dir: Path,
) -> Optional[str]:
    """shape.image.blob 을 images/ 에 저장하고 ref(상대경로) 반환. 실패 시 None."""
    try:
        img = shape.image
        ext = (img.ext or "bin").lstrip(".")
        blob = img.blob
    except Exception:
        return None
    if not blob:
        return None
    _common.mkdir(images_dir)
    filename = f"slide{slide_num:02d}_shape{shape_idx + 1:02d}.{ext}"
    _common.write_bytes(images_dir / filename, blob)
    return f"images/{filename}"


# ====================================================================
# XLSX
# ====================================================================

def _xlsx_clean_nonfinite(src: Path, dst: Path) -> None:
    """xlsx 시트 XML 안 `<v>NaN</v>`/`<v>Infinity</v>`/`<v>-Infinity</v>` 를 제거한 사본을 dst 에 생성.

    원인: 일부 third-party 라이브러리가 만든 xlsx 가 비유한 부동소수점(NaN/Inf) 을 numeric 셀에
    문자열 그대로 기록 → openpyxl 의 `_cast_number → int('NaN')` 에서 ValueError.
    대응: 시트 XML 의 해당 `<v>` 요소만 제거(해당 셀은 빈 셀 처리). 다른 part(images·drawings·
    styles·shared strings) 는 그대로 복사.
    """
    pat = re.compile(rb"<v>(?:NaN|Infinity|-Infinity|INF|-INF)</v>")
    with zipfile.ZipFile(_common.long_str(src), "r") as zin, \
            zipfile.ZipFile(_common.long_str(dst), "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("xl/worksheets/") and item.filename.endswith(".xml"):
                data = pat.sub(b"", data)
            zout.writestr(item, data)


def _safe_load_xlsx_workbooks(
    input_path: Path,
    cleanup_paths: list[Path],
) -> tuple[Any, Any, Path]:
    """openpyxl 로 wb_values(data_only=True) + wb_formulas(data_only=False) 둘 다 로드.

    비표준 셀값(NaN/Infinity 문자열을 numeric 셀에 담은 xlsx) 은 openpyxl 의 strict int cast 로
    ValueError throw. 이 경우 시트 XML 의 비유한값만 제거한 정제본을 임시 파일로 만들어 재시도.
    정제 발생시 임시 파일을 `cleanup_paths` 에 등록(호출자가 finally 에서 unlink).

    반환: (wb_values, wb_formulas, openpyxl_input_path). openpyxl_input_path 는 후속 openpyxl
    호출(이미지 추출 등) 이 같은 정제본을 재사용하도록 path 노출. 정제 불필요시 input_path 그대로.
    """
    import tempfile
    from openpyxl import load_workbook

    def _is_nonfinite_error(e: BaseException) -> bool:
        cur: Optional[BaseException] = e
        while cur is not None:
            msg = str(cur)
            if "NaN" in msg or "Infinity" in msg:
                return True
            cur = cur.__cause__
        return False

    src_str = _common.long_str(input_path)
    try:
        wb_values = load_workbook(src_str, data_only=True)
        wb_formulas = load_workbook(src_str, data_only=False)
        return wb_values, wb_formulas, input_path
    except ValueError as e:
        if not _is_nonfinite_error(e):
            raise

    base = _common._ensure_tmp_base()
    fd, tmp_str = tempfile.mkstemp(prefix="sd-unpack-xlsx-clean-", suffix=".xlsx", dir=str(base))
    os.close(fd)
    cleaned = Path(tmp_str)
    # 등록을 정제·로드 전에 수행 → 도중 throw 해도 호출자 finally 가 unlink.
    cleanup_paths.append(cleaned)
    _xlsx_clean_nonfinite(input_path, cleaned)
    wb_values = load_workbook(_common.long_str(cleaned), data_only=True)
    wb_formulas = load_workbook(_common.long_str(cleaned), data_only=False)
    return wb_values, wb_formulas, cleaned


def _run_xlsx(
    input_path: Path,
    out_dir: Path,
    *,
    source_name_override: Optional[str] = None,
    tool_extra: str = "",
) -> None:
    _common.ensure_pip("openpyxl")
    from openpyxl import load_workbook
    from openpyxl.worksheet.worksheet import Worksheet

    sheets_dir = out_dir / "sheets"
    charts_dir = out_dir / "charts"
    sheet_summaries: list[str] = []
    sheet_names: list[tuple[str, str, str]] = []  # (idx, safe_name, raw_name) — 일반 Worksheet
    chart_sheet_names: list[tuple[str, str, str]] = []  # (idx, safe_name, raw_name) — Chartsheet
    sheet_charts: dict[str, list[str]] = {}  # idx -> chart filenames
    sheet_formula_count: dict[str, int] = {}
    sheet_dims: dict[str, tuple[int, int]] = {}

    # 비표준 셀값(NaN/Infinity) 사전 정제 + openpyxl 로드. 정제본 임시파일은 마지막에 unlink.
    _xlsx_cleanups: list[Path] = []
    try:
        wb_values, wb_formulas, openpyxl_input = _safe_load_xlsx_workbooks(input_path, _xlsx_cleanups)
        try:
            _common.mkdir(sheets_dir)
            # openpyxl 의 sheetnames 는 일반 Worksheet 와 Chartsheet 둘 다 포함.
            # 시트 순서 그대로 idx 통합 부여 (사용자 워크북 순서 보존).
            # 일반 Worksheet 만 COM Excel PNG export 대상, Chartsheet 는 차트 데이터만 추출.
            idx_counter = 0
            for sheet_name in wb_values.sheetnames:
                obj = wb_values[sheet_name]
                idx_counter += 1
                idx = f"{idx_counter:02d}"
                safe_name = _common.slugify_filename(sheet_name, max_len=40)
                if isinstance(obj, Worksheet):
                    sheet_names.append((idx, safe_name, sheet_name))
                else:
                    # Chartsheet 등 비-worksheet
                    chart_sheet_names.append((idx, safe_name, sheet_name))

            # COM Excel 호출: 데이터 영역 → ChartObject + Range.CopyPicture → 시트별 PNG.
            # 시트별 (last_row, last_col) 도 같이 반환되어 .jsonl 이 같은 데이터 영역으로 통일됨.
            # PNG export 실패한 시트는 sheet_png_skipped 에 사유 (silent skip 금지).
            with _common.com_lock():
                # openpyxl_input 사용: 정제본(NaN 제거) 이 있으면 COM Excel 도 정제본을 열어야 함
                # (Excel 역시 `<v>NaN</v>` 가 있는 xlsx 의 Open 에 실패).
                sheet_ranges, sheet_png_skipped = _excel_export_sheet_pngs(openpyxl_input, sheets_dir, sheet_names)

            for idx, safe_name, raw_name in sheet_names:
                ws_v = wb_values[raw_name]
                ws_f = wb_formulas[raw_name]

                # COM Find 결과가 있으면 그 범위, 없으면 openpyxl max_row/max_column fallback.
                last_row, last_col = sheet_ranges.get(raw_name, (ws_v.max_row, ws_v.max_column))
                sheet_dims[idx] = (last_row, last_col)

                jsonl_lines, formula_n = _sheet_to_jsonl(ws_v, ws_f, last_row, last_col)
                _common.write_text(sheets_dir / f"{idx}_{safe_name}.jsonl", "\n".join(jsonl_lines))
                sheet_formula_count[idx] = formula_n

                for chart_idx, chart in enumerate(getattr(ws_f, "_charts", []), start=1):
                    data = _extract_openpyxl_chart_data(chart)
                    _common.mkdir(charts_dir)
                    chart_filename = f"sheet{idx}_chart{chart_idx:02d}.data.json"
                    _common.write_text(
                        charts_dir / chart_filename,
                        json.dumps(data, ensure_ascii=False, indent=2),
                    )
                    sheet_charts.setdefault(idx, []).append(chart_filename)

            # Chartsheet 처리: 차트 데이터를 charts/sheet<idx>_chart.data.json 으로 저장
            chart_sheet_chart_files: dict[str, str] = {}  # idx -> chart filename
            for idx, safe_name, raw_name in chart_sheet_names:
                cs = wb_formulas[raw_name]
                chart = None
                # Chartsheet.charts 또는 _charts 속성 (openpyxl 버전 따라 다름)
                for attr in ("charts", "_charts"):
                    v = getattr(cs, attr, None)
                    if v:
                        if hasattr(v, "__iter__"):
                            try:
                                chart = next(iter(v), None)
                            except Exception:
                                chart = None
                        else:
                            chart = v
                        if chart is not None:
                            break
                if chart is None:
                    # 단일 chart 속성 fallback
                    chart = getattr(cs, "chart", None)
                if chart is not None:
                    try:
                        data = _extract_openpyxl_chart_data(chart)
                    except Exception:
                        data = None
                    if data is not None:
                        _common.mkdir(charts_dir)
                        chart_filename = f"sheet{idx}_chart.data.json"
                        _common.write_text(
                            charts_dir / chart_filename,
                            json.dumps(data, ensure_ascii=False, indent=2),
                        )
                        chart_sheet_chart_files[idx] = chart_filename

            # 워크북 단위 메타 (defined names·pivots·sheet codeName 등) — 시트 jsonl 외부 분리.
            wb_meta = _workbook_meta(wb_formulas, input_path)
            # VBA 시트 객체명 ↔ raw 시트명 매핑 (시트 codeName 기반)
            sheet_code_map: dict[str, str] = {}
            for ws in wb_formulas.worksheets:
                code = getattr(ws.sheet_properties, "codeName", None)
                if code:
                    sheet_code_map[code] = ws.title
            if sheet_code_map:
                wb_meta["sheet_code_map"] = sheet_code_map
            if wb_meta:
                _common.write_text(
                    out_dir / "workbook.meta.json",
                    json.dumps(wb_meta, ensure_ascii=False, indent=2),
                )
        finally:
            wb_values.close()
            wb_formulas.close()

        # 시트 PNG 는 데이터 영역(Find 범위) 만 캡처 → 데이터 영역 밖 이미지는 누락될 수 있음 →
        # raw 이미지를 시트+셀 위치 정보 포함해서 별도 보존.
        # openpyxl_input 사용: 정제본이 있으면 같은 정제본으로 로드(원본은 openpyxl 가 못 읽음).
        sheet_images = _extract_xlsx_images_with_position(openpyxl_input, out_dir, sheet_names)
        attachment_links = _extract_zip_media(
            input_path,
            out_dir,
            media_zip_prefix="xl/media/",
            embed_zip_prefix="xl/embeddings/",
        )

        # 시트별 산출물 풀목록 — 일반 시트 + chart sheet 통합, 시트 순서 (idx) 대로
        sheet_summary_map: dict[str, str] = {}
        for idx, safe_name, raw_name in sheet_names:
            last_row, last_col = sheet_dims.get(idx, (0, 0))
            formula_n = sheet_formula_count.get(idx, 0)
            png_path = sheets_dir / f"{idx}_{safe_name}.png"
            if png_path.exists():
                parts = [f"`sheets/{idx}_{safe_name}.png`", "`.jsonl`"]
            else:
                # PNG 미생성 — worker 가 사유 전달 (16-bit cap / COM 실패 등)
                reason = sheet_png_skipped.get(raw_name, "사유 미상")
                parts = [f"`sheets/{idx}_{safe_name}.jsonl`", f"(PNG 미생성 — {reason})"]
            chart_refs = sheet_charts.get(idx, [])
            if chart_refs:
                parts.append("(차트: " + ", ".join(f"`charts/{c}`" for c in chart_refs) + ")")
            img_refs = sheet_images.get(raw_name, [])
            if img_refs:
                parts.append("(이미지: " + ", ".join(f"`images/{n}`" for n in img_refs) + ")")
            meta = f"({last_row}행×{last_col}열"
            if formula_n:
                meta += f", 수식 {formula_n}개"
            meta += ")"
            sheet_summary_map[idx] = " ".join(parts) + " " + meta

        for idx, safe_name, raw_name in chart_sheet_names:
            chart_filename = chart_sheet_chart_files.get(idx)
            if chart_filename:
                sheet_summary_map[idx] = f"`charts/{chart_filename}` (chart sheet — \"{raw_name}\")"
            else:
                sheet_summary_map[idx] = f"(chart sheet — \"{raw_name}\", 차트 데이터 추출 실패)"

        # idx 순서대로 통합
        for idx in sorted(sheet_summary_map.keys()):
            sheet_summaries.append(sheet_summary_map[idx])

        source_name, source_size = _source_meta(input_path, out_dir, source_name_override)
        macro_modules = _extract_macros(_source_path(out_dir, source_name), out_dir)

        sections: dict[str, list[str]] = {}
        if sheet_summaries:
            sections[f"시트 (총 {len(sheet_summaries)}개)"] = sheet_summaries
        if macro_modules:
            sections[f"VBA 매크로 (총 {len(macro_modules)}개)"] = [f"`macros/{m}`" for m in macro_modules]

        _common.write_readme(
            out_dir,
            source_name=source_name,
            source_size=source_size,
            tool=("openpyxl + COM Excel + ZIP " + tool_extra).strip(),
            loss_notes=(
                "셀 서식(바탕색·border·폰트)·frozen·dims 미보존 (필요 시 _source.xlsx 직접 추출). "
                "시각은 시트별 PNG, 분석 데이터(셀값·number_format·수식·merges·hyperlinks·comments) 는 "
                "시트별 .jsonl 한 줄=한 행(좌표 명시), 워크북 단위 메타(defined names 등) 는 workbook.meta.json."
            ),
            sections=sections or None,
            attachments=attachment_links,
        )
    finally:
        for _p in _xlsx_cleanups:
            try:
                _p.unlink()
            except Exception:
                pass


# ====================================================================
# COM 헬퍼 — 모두 office_worker.py subprocess 로 격리.
# 이유: COM IUnknown::Release RPC release 가 Office 프로세스 cleanup 비동기성 때문에 60s 까지 wait 가능.
# 별도 process 로 격리하면 process 종료 = OS 가 COM cleanup 강제 → wait 없음 + hang 시 timeout kill.
# ====================================================================

_WORKER_PATH = Path(__file__).parent / "office_worker.py"


def _run_worker(cmd: str, *args: str, timeout: float, capture_stdout: bool = False) -> str:
    """office_worker.py subprocess 호출. timeout 초과 시 kill 후 throw.

    호출자는 com_lock 으로 sequential 보장. worker 자체는 단독 process 라 인스턴스 격리됨.
    """
    # worker 는 동일 sys.executable 로 실행되므로 호출 전 호스트 환경에 pywin32 보장 → worker 가 import 가능.
    _common.ensure_pip("pythoncom", "pywin32")
    import subprocess
    # errors="replace": Windows COM 에러 메시지가 cp949 등 비-utf8 로 올 수 있음 → None 변환 회피.
    proc = subprocess.run(
        [sys.executable, str(_WORKER_PATH), cmd, *args],
        capture_output=True, text=True, timeout=timeout, encoding="utf-8", errors="replace",
    )
    if proc.returncode != 0:
        raise RuntimeError(
            f"office_worker '{cmd}' failed (rc={proc.returncode}): {proc.stderr.strip()}"
        )
    return proc.stdout if capture_stdout else ""


def _word_export_pdf(input_path: Path, output_pdf: Path) -> None:
    _run_worker("word_pdf", str(input_path), str(output_pdf), timeout=600)


def _powerpoint_export_slides(
    input_path: Path,
    out_dir: Path,
    slide_titles: list[tuple[str, str]],
) -> None:
    """슬라이드별 PNG 직접 export. worker subprocess 로 격리."""
    _run_worker(
        "ppt_png", str(input_path), str(out_dir), json.dumps(slide_titles),
        timeout=600,
    )


def _excel_export_sheet_pngs(
    input_path: Path,
    sheets_dir: Path,
    sheet_names: list[tuple[str, str, str]],
) -> tuple[dict[str, tuple[int, int]], dict[str, str]]:
    """시트별 PNG 생성 + (last_row, last_col) 매핑 + skipped 사유 반환.

    호출자에서 sheetProtection strip 사본 만들고 worker 에 그 사본 path 만 넘김.
    Excel COM 자체 작업은 worker subprocess.

    반환: (sheet_ranges, skipped) — skipped 는 PNG export 실패한 시트의 사유 dict (raw_name → reason).
    """
    with _common.temp_workdir() as tmp:
        unprotected = tmp / "_unprotected.xlsx"
        _xlsx_strip_protection(input_path, unprotected)
        result = _run_worker(
            "excel_sheets", str(unprotected), str(sheets_dir), json.dumps(sheet_names),
            timeout=600, capture_stdout=True,
        )
    if not result.strip():
        return {}, {}
    parsed = json.loads(result)
    ranges_raw = parsed.get("sheet_ranges", {})
    sheet_ranges = {k: tuple(v) for k, v in ranges_raw.items()}
    skipped = parsed.get("skipped", {})
    return sheet_ranges, skipped


def _xlsx_strip_protection(src: Path, dst: Path) -> None:
    """xlsx ZIP 안 sheetProtection / workbookProtection xml 노드 제거한 사본 생성.

    Excel COM 의 ChartObjects.Add 가 보호 시트에서 차단되는 이슈 회피용.
    비번 hash 우회 불필요 — xlsx 의 보호는 xml 노드가 본체. 다른 part 는 그대로 → 차트/매크로/이미지 보존.
    """
    import re
    # attribute value 에 '/' 들어갈 수 있음 (hashValue base64) → '/' 제외하면 안 됨.
    # '>' 만 제외하고 lazy 로 첫 '/>' 까지.
    sheet_re = re.compile(rb'<sheetProtection\b[^>]*?/>')
    wb_re = re.compile(rb'<workbookProtection\b[^>]*?/>')
    with zipfile.ZipFile(_common.long_str(src), "r") as zin, \
         zipfile.ZipFile(_common.long_str(dst), "w", zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info)
            entry_name = info.filename
            if entry_name.startswith("xl/worksheets/sheet") and entry_name.endswith(".xml"):
                data = sheet_re.sub(b"", data)
            elif entry_name == "xl/workbook.xml":
                data = wb_re.sub(b"", data)
            zout.writestr(info, data)


def _convert_legacy(input_path: Path, target_path: Path) -> None:
    """input_path (.doc/.ppt/.xls/.xlsb) 를 신형 형식으로 변환해 target_path 에 저장.

    임시 폴더에 변환 후 long-path-safe copy. Office COM 작업은 worker subprocess.
    """
    target_ext = target_path.suffix
    with _common.com_lock(), _common.temp_workdir() as tmp:
        # Excel SaveAs 등 Office COM 은 파일명에 `< > ? [ ] : | * "` 를 거부 (대괄호는
        # 워크북 참조 구문으로 해석). 임시 변환 파일명에서만 치환 — 최종 결과는
        # target_path 의 원본 이름으로 copy 되므로 출력명에는 영향 없음.
        safe_stem = re.sub(r'[<>?\[\]:|*"]', "_", input_path.stem)
        tmp_target = tmp / (safe_stem + target_ext)
        _run_worker(
            "convert_legacy", input_path.suffix.lower(), str(input_path), str(tmp_target),
            timeout=600,
        )
        _common.copy(tmp_target, target_path)


# ====================================================================
# 공용
# ====================================================================

def _render_pdf_pages(pdf_path: Path, pages_dir: Path) -> list[str]:
    """PDF → 페이지별 PNG + MD."""
    _common.ensure_pip("fitz", "PyMuPDF")
    import fitz

    summaries: list[str] = []
    doc = fitz.open(_common.long_str(pdf_path))
    try:
        for i, page in enumerate(doc, start=1):
            idx = f"{i:03d}"
            text = page.get_text("text") or ""
            _common.write_text(pages_dir / f"{idx}.md", text)
            pix = page.get_pixmap(dpi=300)
            pix.save(_common.long_str(pages_dir / f"{idx}.png"))
            summaries.append(f"`pages/{idx}.png` (시각) — `.md` ({len(text)}자)")
    finally:
        doc.close()
    return summaries


def _source_meta(
    input_path: Path,
    out_dir: Path,
    source_name_override: Optional[str],
) -> tuple[str, int]:
    if source_name_override:
        src_ext = Path(source_name_override).suffix.lstrip(".")
        src_path = out_dir / f"_source.{src_ext}"
        return source_name_override, src_path.stat().st_size
    return input_path.name, input_path.stat().st_size


def _source_path(out_dir: Path, source_name: str) -> Path:
    """out_dir 안 보존된 원본 _source.<ext> 경로."""
    ext = Path(source_name).suffix.lstrip(".")
    return out_dir / f"_source.{ext}"


def _extract_macros(input_path: Path, out_dir: Path) -> list[str]:
    """OLE/OOXML 파일에서 VBA 매크로 추출. macros/<모듈명>.vba 로 저장 (원본 코드 그대로).

    추출된 모듈 파일명 list 반환 (예: ["Module1.vba", "ThisWorkbook.vba"]).
    매크로 없으면 빈 list.

    시트 객체명↔raw 시트명 매핑은 호출자(_run_xlsx)가 workbook.meta.json 에 별도 보관.
    """
    _common.ensure_pip("oletools")
    from oletools.olevba import VBA_Parser

    module_files: list[str] = []
    parser = VBA_Parser(_common.long_str(input_path))
    try:
        if not parser.detect_vba_macros():
            return module_files
        macros_dir = out_dir / "macros"
        _common.mkdir(macros_dir)
        for (_filename, stream_path, vba_filename, vba_code) in parser.extract_macros():
            module_name = vba_filename or stream_path or "module"
            stem = Path(module_name).stem or "module"
            dst = _common.unique_path(macros_dir, f"{stem}.vba")
            _common.write_text(dst, vba_code or "")
            module_files.append(dst.name)
        return module_files
    finally:
        parser.close()


def _extract_xlsx_images_with_position(
    input_path: Path,
    out_dir: Path,
    sheet_names: list[tuple[str, str, str]],
) -> dict[str, list[str]]:
    """xlsx 의 이미지를 시트+셀 위치 정보 포함 파일명으로 images/ 에 저장.

    파일명 패턴: `<sheet_safe_name>_<cell_addr>.<ext>` (예: `BOA_E5.png`).
    같은 셀에 여러 이미지면 unique_path 가 _1, _2 suffix.
    이미지 데이터/anchor 추출 실패 시 그 이미지만 skip (전체 fail 안 함).
    """
    _common.ensure_pip("openpyxl")
    from openpyxl import load_workbook
    from openpyxl.utils import get_column_letter

    images_dir = out_dir / "images"
    name_to_safe = {raw: safe for _idx, safe, raw in sheet_names}
    sheet_images: dict[str, list[str]] = {}  # raw_name -> [filename, ...]

    wb = load_workbook(_common.long_str(input_path))
    try:
        for sheet_name in wb.sheetnames:
            if sheet_name not in name_to_safe:
                continue  # chart sheet 등 일반 worksheet 아닌 경우 skip
            ws = wb[sheet_name]
            safe_name = name_to_safe[sheet_name]
            for img in getattr(ws, "_images", []):
                # anchor 위치 추출 (OneCellAnchor / TwoCellAnchor 둘 다 _from 갖음)
                cell_addr = "anchor_unknown"
                anchor = img.anchor
                marker = getattr(anchor, "_from", None) or getattr(anchor, "from", None)
                if marker is not None:
                    col = marker.col + 1  # 0-based → 1-based
                    row = marker.row + 1
                    cell_addr = f"{get_column_letter(col)}{row}"

                # 이미지 raw 데이터 추출 (openpyxl 내부 API).
                data = img._data()
                if not data:
                    continue

                # 확장자: magic bytes 로 추정.
                ext = "bin"
                if data.startswith(b"\x89PNG"):
                    ext = "png"
                elif data.startswith(b"\xff\xd8"):
                    ext = "jpg"
                elif data.startswith(b"GIF8"):
                    ext = "gif"
                elif data.startswith(b"BM"):
                    ext = "bmp"
                elif data[:4] == b"RIFF" and data[8:12] == b"WEBP":
                    ext = "webp"

                _common.mkdir(images_dir)
                base = f"{safe_name}_{cell_addr}.{ext}"
                dst = _common.unique_path(images_dir, base)
                _common.write_bytes(dst, data)
                sheet_images.setdefault(sheet_name, []).append(dst.name)
    finally:
        wb.close()
    return sheet_images


def _extract_zip_media(
    input_path: Path,
    out_dir: Path,
    *,
    media_zip_prefix: str,
    embed_zip_prefix: str,
    images_dir: Optional[Path] = None,
) -> list[str]:
    """OOXML(docx/pptx/xlsx) ZIP 안의 media/ 와 embeddings/ 에서 이미지/첨부 추출.

    images_dir=None 이면 media/ 추출 skip (PDF 가 시각 책임지는 형식: docx, pptx).
    xlsx 는 별도 _extract_xlsx_images_with_position 으로 위치 정보 포함 추출.
    """
    attachment_links: list[str] = []
    attachments_dir = out_dir / "attachments"

    with zipfile.ZipFile(_common.long_str(input_path), "r") as zf:
        for info in zf.infolist():
            entry_name = info.filename
            if info.is_dir():
                continue
            if images_dir is not None and entry_name.startswith(media_zip_prefix):
                _common.mkdir(images_dir)
                base = Path(entry_name).name
                dst = _common.unique_path(images_dir, base)
                with zf.open(info) as f:
                    _common.write_bytes(dst, f.read())
            elif entry_name.startswith(embed_zip_prefix):
                _common.mkdir(attachments_dir)
                base = Path(entry_name).name
                dst = _common.unique_path(attachments_dir, base)
                with zf.open(info) as f:
                    _common.write_bytes(dst, f.read())
                size = dst.stat().st_size
                recursed = maybe_recurse_attachment(dst, attachments_dir)
                if recursed is not None:
                    os.unlink(_common.long_str(dst))
                    attachment_links.append(f"attachments/{recursed.name}/ ({_common.format_size(size)})")
                else:
                    attachment_links.append(f"attachments/{dst.name} ({_common.format_size(size)})")
    return attachment_links


def _json_default(obj: Any) -> str:
    """JSON 직렬화 fallback. openpyxl datetime → ISO 8601. 그 외는 throw."""
    if isinstance(obj, (datetime, date, time)):
        return obj.isoformat()
    raise TypeError(f"not JSON serializable: {type(obj).__name__}")


def _sheet_to_jsonl(
    ws_v, ws_f, last_row: int, last_col: int,
) -> tuple[list[str], int]:
    """openpyxl Worksheet 의 (1,1)~(last_row,last_col) → 행 단위 JSONL.

    분석 핵심: 데이터·number_format·수식. 시각 표시(바탕색·border·폰트·frozen)·dims 는 미보존
    (필요 시 Claude 가 _source.xlsx 직접 추출).

    데이터 jsonl (한 줄=한 행. 빈 셀 키 생략):
    - 첫 줄: `{"_meta":{"merges":[...], "number_formats":{...}, "hyperlinks":{...}, "comments":{...}}}`
      - merges: 셀 좌표 해석 필수 (머지 영역 안 빈 셀 오해 차단)
      - number_formats: Date·통화·% 등 셀 값 의미 단서
      - hyperlinks·comments: 셀 부가 정보
      - 비어있는 키는 생략
    - 데이터 줄: `{"r":N, "<col>":value, ..., "_f":{<col>:formula}}`
    - 빈 행도 `{"r":N}` 한 줄 유지

    반환: (lines, formula_count)
    """
    from openpyxl.utils import get_column_letter

    if last_row < 1 or last_col < 1:
        return [json.dumps({"_meta": {}}, ensure_ascii=False)], 0

    meta: dict[str, Any] = {}
    merges = [str(r) for r in ws_v.merged_cells.ranges]
    if merges:
        meta["merges"] = merges

    hyperlinks: dict[str, str] = {}
    comments: dict[str, str] = {}
    number_formats: dict[str, str] = {}
    for row in ws_v.iter_rows(min_row=1, max_row=last_row, min_col=1, max_col=last_col):
        for cell in row:
            hl = getattr(cell, "hyperlink", None)
            if hl is not None and getattr(hl, "target", None):
                hyperlinks[cell.coordinate] = hl.target
            cm = getattr(cell, "comment", None)
            if cm is not None and getattr(cm, "text", None):
                comments[cell.coordinate] = cm.text
            nf = getattr(cell, "number_format", None)
            if nf and nf != "General":
                number_formats[cell.coordinate] = nf
    if number_formats:
        meta["number_formats"] = number_formats
    if hyperlinks:
        meta["hyperlinks"] = hyperlinks
    if comments:
        meta["comments"] = comments

    lines: list[str] = [json.dumps({"_meta": meta}, ensure_ascii=False, default=_json_default)]
    formula_count = 0

    rows_v = ws_v.iter_rows(min_row=1, max_row=last_row, min_col=1, max_col=last_col, values_only=True)
    rows_f = ws_f.iter_rows(min_row=1, max_row=last_row, min_col=1, max_col=last_col)
    for r_idx, (row_v, row_f) in enumerate(zip(rows_v, rows_f), start=1):
        row_data: dict[str, Any] = {"r": r_idx}
        fmap: dict[str, str] = {}
        for c_idx, (v, fcell) in enumerate(zip(row_v, row_f), start=1):
            col_letter = get_column_letter(c_idx)
            if v is not None:
                row_data[col_letter] = v
            if fcell.data_type == "f":
                fv = fcell.value
                # 일반·shared formula 는 str, array formula 는 ArrayFormula(.text 보유)
                fmap[col_letter] = fv if isinstance(fv, str) else getattr(fv, "text", str(fv))
                formula_count += 1
        if fmap:
            row_data["_f"] = fmap
        lines.append(json.dumps(row_data, ensure_ascii=False, default=_json_default))

    return lines, formula_count


def _workbook_meta(wb, input_path: Path) -> dict[str, Any]:
    """워크북 단위 메타 (defined names·pivot tables 등). 비어있으면 빈 dict 반환."""
    meta: dict[str, Any] = {}
    defined_names: dict[str, list[str]] = {}
    # openpyxl 3.x: wb.defined_names 는 DefinedNameDict (dict-like)
    try:
        for def_name, dn in wb.defined_names.items():
            try:
                dests = [f"'{sheet}'!{addr}" for sheet, addr in dn.destinations]
            except Exception:
                # destinations 파싱 불가 시 raw value 보존 (예: 워크북-수식 형태)
                dests = [str(getattr(dn, "value", ""))]
            defined_names[def_name] = dests
    except Exception:
        # defined_names 자체 접근 실패 → 워크북에 없는 것으로 처리
        pass
    if defined_names:
        meta["defined_names"] = defined_names

    pivots = _extract_pivots(input_path)
    if pivots:
        meta["pivots"] = pivots

    return meta


_XLSX_NS = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"
_XLSX_REL_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
_PKG_REL_NS = "{http://schemas.openxmlformats.org/package/2006/relationships}"


def _extract_pivots(input_path: Path) -> list[dict]:
    """xlsx 의 pivot table 정의 list. ZIP 안 `xl/pivotTables/*.xml` + `xl/pivotCache/*.xml` 파싱.

    cacheId 매핑은 workbook.xml 의 pivotCaches + workbook.xml.rels 통해 정확히 해결.
    - workbook.xml 의 pivotCaches: cacheId → r:id
    - workbook.xml.rels: Id → Target (cache xml 파일)

    각 pivot 의 정보:
    - name: pivot table 이름
    - location: 펼쳐진 위치 (예: "A1:E20")
    - source: 원본 데이터 위치 (예: "'Sheet1'!A1:D100")
    - rowFields/colFields/pageFields: 행·열·필터 필드명 list
    - dataFields: 값 필드 [{name, field, subtotal}, ...] (subtotal = sum/count/average/...)
    """
    import xml.etree.ElementTree as ET

    pivots: list[dict] = []
    try:
        with zipfile.ZipFile(_common.long_str(input_path), "r") as zf:
            namelist = zf.namelist()
            pivot_files = sorted(
                n for n in namelist
                if n.startswith("xl/pivotTables/pivotTable") and n.endswith(".xml")
            )
            if not pivot_files:
                return pivots

            # 1. workbook.xml.rels 에서 Id → Target 매핑
            rid_to_target: dict[str, str] = {}
            try:
                rels_root = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))
                for rel in rels_root.findall(f"{_PKG_REL_NS}Relationship"):
                    rid_to_target[rel.get("Id", "")] = rel.get("Target", "")
            except Exception:
                pass

            # 2. workbook.xml 의 pivotCaches 에서 cacheId → cache 파일 경로 매핑
            cache_id_to_file: dict[str, str] = {}
            try:
                wb_root = ET.fromstring(zf.read("xl/workbook.xml"))
                pcs = wb_root.find(f"{_XLSX_NS}pivotCaches")
                if pcs is not None:
                    for pc in pcs:
                        cid = pc.get("cacheId")
                        rid = pc.get(f"{_XLSX_REL_NS}id")
                        if not cid or not rid:
                            continue
                        target = rid_to_target.get(rid, "")
                        if not target:
                            continue
                        # target 의 상대 경로 → ZIP 안 절대 경로
                        if target.startswith("/"):
                            cache_path = target.lstrip("/")
                        else:
                            cache_path = "xl/" + target
                        cache_id_to_file[cid] = cache_path
            except Exception:
                pass

            # 3. cache 파일 파싱: cacheId → {source, field_names}
            cache_info: dict[str, dict] = {}
            for cid, cf in cache_id_to_file.items():
                try:
                    root = ET.fromstring(zf.read(cf))
                except Exception:
                    continue
                info: dict = {}
                cs = root.find(f"{_XLSX_NS}cacheSource")
                if cs is not None:
                    ws = cs.find(f"{_XLSX_NS}worksheetSource")
                    if ws is not None:
                        sheet = ws.get("sheet", "")
                        ref = ws.get("ref", "")
                        named = ws.get("name", "")
                        if sheet and ref:
                            info["source"] = f"'{sheet}'!{ref}"
                        elif named:
                            info["source"] = named
                fields_elem = root.find(f"{_XLSX_NS}cacheFields")
                if fields_elem is not None:
                    field_names: list[str] = []
                    for f in fields_elem:
                        if f.tag == f"{_XLSX_NS}cacheField":
                            field_names.append(f.get("name", ""))
                    info["field_names"] = field_names
                cache_info[cid] = info

            # pivot table 파일 파싱
            for pf in pivot_files:
                try:
                    root = ET.fromstring(zf.read(pf))
                except Exception:
                    continue
                pivot: dict = {"name": root.get("name", "")}
                cache_id = root.get("cacheId", "")
                field_names: list[str] = []
                if cache_id and cache_id in cache_info:
                    ci = cache_info[cache_id]
                    if "source" in ci:
                        pivot["source"] = ci["source"]
                    field_names = ci.get("field_names", [])

                loc = root.find(f"{_XLSX_NS}location")
                if loc is not None:
                    pivot["location"] = loc.get("ref", "")

                # row·col·page fields (인덱스 → 이름)
                for tag, key in (
                    ("rowFields", "rowFields"),
                    ("colFields", "colFields"),
                    ("pageFields", "pageFields"),
                ):
                    elem = root.find(f"{_XLSX_NS}{tag}")
                    if elem is None:
                        continue
                    names: list[str] = []
                    for child in elem:
                        x = child.get("x") or child.get("fld")
                        if x is None:
                            continue
                        try:
                            idx = int(x)
                        except (TypeError, ValueError):
                            continue
                        if 0 <= idx < len(field_names) and field_names[idx]:
                            names.append(field_names[idx])
                        else:
                            names.append(f"field_{idx}")
                    if names:
                        pivot[key] = names

                # dataFields (값 필드 + 집계 함수)
                df_elem = root.find(f"{_XLSX_NS}dataFields")
                if df_elem is not None:
                    df_list: list[dict] = []
                    for df in df_elem:
                        if df.tag != f"{_XLSX_NS}dataField":
                            continue
                        fld = df.get("fld", "")
                        field_name = ""
                        try:
                            idx = int(fld)
                            if 0 <= idx < len(field_names):
                                field_name = field_names[idx]
                        except (TypeError, ValueError):
                            pass
                        df_list.append({
                            "name": df.get("name", ""),
                            "field": field_name,
                            "subtotal": df.get("subtotal", "sum"),
                        })
                    if df_list:
                        pivot["dataFields"] = df_list

                pivots.append(pivot)
    except Exception as e:
        raise RuntimeError(f"pivot 추출 실패 ({input_path.name}) — {e}") from e
    return pivots


def _extract_pptx_chart_data(chart) -> dict:
    data: dict = {"type": str(getattr(chart, "chart_type", "")), "categories": [], "series": []}
    if chart.plots:
        cats = list(chart.plots[0].categories)
        data["categories"] = [str(c) for c in cats]
    for plot in chart.plots:
        for series in plot.series:
            values = list(series.values) if hasattr(series, "values") else []
            data["series"].append({
                "name": str(getattr(series, "name", "")),
                "values": [v for v in values],
            })
    return data


def _extract_openpyxl_chart_data(chart) -> dict:
    data: dict = {"type": type(chart).__name__, "title": "", "series": []}
    if chart.title:
        data["title"] = str(chart.title)
    for s in chart.series:
        ref = ""
        if s.val and s.val.numRef:
            ref = str(s.val.numRef.f)
        data["series"].append({"ref": ref})
    return data
