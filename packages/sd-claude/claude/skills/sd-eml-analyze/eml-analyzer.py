#!/usr/bin/env python3
"""EML Email Analyzer - Parses EML files and attachments into structured markdown."""

import sys
import os
import io
import subprocess
import email
import html
import re
import tempfile
from email.policy import default as default_policy
from email.header import decode_header
from pathlib import Path

# stdout UTF-8 강제 (Windows 호환)
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8", errors="replace")


def ensure_packages():
    """필요한 패키지 자동 설치."""
    packages = {
        "pdfminer.six": "pdfminer",
        "python-pptx": "pptx",
        "openpyxl": "openpyxl",
    }
    missing = []
    for pip_name, import_name in packages.items():
        try:
            __import__(import_name)
        except ImportError:
            missing.append(pip_name)
    if missing:
        print(f"패키지 설치 중: {', '.join(missing)}...", file=sys.stderr)
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "-q", *missing],
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )


ensure_packages()

from pdfminer.high_level import extract_text as pdf_extract_text  # noqa: E402
from pptx import Presentation  # noqa: E402
from openpyxl import load_workbook  # noqa: E402


# ── Korean charset helpers ──────────────────────────────────────────

KOREAN_CHARSET_MAP = {
    "ks_c_5601-1987": "cp949",
    "ks_c_5601": "cp949",
    "euc_kr": "cp949",
    "euc-kr": "cp949",
}


def fix_charset(charset: str) -> str:
    if charset is None:
        return "utf-8"
    return KOREAN_CHARSET_MAP.get(charset.lower(), charset)


# ── EML parsing ─────────────────────────────────────────────────────


def parse_eml(filepath: str):
    with open(filepath, "rb") as f:
        msg = email.message_from_binary_file(f, policy=default_policy)

    # Headers
    headers = {
        "subject": str(msg["Subject"] or ""),
        "from": str(msg["From"] or ""),
        "to": str(msg["To"] or ""),
        "cc": str(msg["Cc"] or ""),
        "date": str(msg["Date"] or ""),
    }

    # Body
    body_plain = ""
    body_html = ""

    if msg.is_multipart():
        for part in msg.walk():
            ctype = part.get_content_type()
            cdisp = part.get_content_disposition()
            if cdisp == "attachment":
                continue
            if ctype == "text/plain" and not body_plain:
                body_plain = _get_text(part)
            elif ctype == "text/html" and not body_html:
                body_html = _get_text(part)
    else:
        body_plain = _get_text(msg)

    # Attachments
    attachments = []
    for part in msg.walk():
        filename = part.get_filename()
        if not filename:
            continue
        cdisp = part.get_content_disposition()
        if cdisp not in ("attachment", "inline", None):
            continue
        payload = part.get_payload(decode=True)
        if payload is None:
            continue
        attachments.append(
            {
                "filename": filename,
                "content_type": part.get_content_type(),
                "size": len(payload),
                "data": payload,
            }
        )

    return headers, body_plain, body_html, attachments


def _get_text(part) -> str:
    try:
        return part.get_content()
    except Exception:
        payload = part.get_payload(decode=True)
        if not payload:
            return ""
        charset = fix_charset(part.get_content_charset())
        return payload.decode(charset, errors="replace")


# ── Attachment extractors ───────────────────────────────────────────


def extract_pdf(data: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
        f.write(data)
        tmp = f.name
    try:
        text = pdf_extract_text(tmp)
        return text.strip() if text else "(텍스트 추출 실패)"
    except Exception as e:
        return f"(PDF 파싱 오류: {e})"
    finally:
        os.unlink(tmp)


def extract_pptx(data: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        f.write(data)
        tmp = f.name
    try:
        prs = Presentation(tmp)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            lines = [f"#### 슬라이드 {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs)
                        if line.strip():
                            lines.append(line)
                if shape.has_table:
                    header = " | ".join(
                        cell.text for cell in shape.table.rows[0].cells
                    )
                    sep = " | ".join(
                        "---" for _ in shape.table.rows[0].cells
                    )
                    lines.append(f"| {header} |")
                    lines.append(f"| {sep} |")
                    for row in list(shape.table.rows)[1:]:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        lines.append(f"| {row_text} |")
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text
                if notes.strip():
                    lines.append(f"> 노트: {notes}")
            slides.append("\n".join(lines))
        return "\n\n".join(slides) if slides else "(텍스트 없음)"
    except Exception as e:
        return f"(PPTX 파싱 오류: {e})"
    finally:
        os.unlink(tmp)


def extract_xlsx(data: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as f:
        f.write(data)
        tmp = f.name
    try:
        wb = load_workbook(tmp, data_only=True)
        sheets = []
        for name in wb.sheetnames:
            ws = wb[name]
            lines = [f"#### 시트: {name}"]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                lines.append("(데이터 없음)")
                sheets.append("\n".join(lines))
                continue
            # 마크다운 테이블
            first_row = rows[0]
            col_count = len(first_row)
            header = " | ".join(str(c) if c is not None else "" for c in first_row)
            sep = " | ".join("---" for _ in range(col_count))
            lines.append(f"| {header} |")
            lines.append(f"| {sep} |")
            for row in rows[1:]:
                vals = " | ".join(str(c) if c is not None else "" for c in row)
                if any(c is not None for c in row):
                    lines.append(f"| {vals} |")
            sheets.append("\n".join(lines))
        return "\n\n".join(sheets) if sheets else "(데이터 없음)"
    except Exception as e:
        return f"(XLSX 파싱 오류: {e})"
    finally:
        os.unlink(tmp)


def extract_text_file(data: bytes) -> str:
    for enc in ("utf-8", "cp949", "euc-kr", "latin-1"):
        try:
            return data.decode(enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return data.decode("utf-8", errors="replace")


# ── HTML stripping ──────────────────────────────────────────────────


def strip_html(text: str) -> str:
    text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.I)
    text = re.sub(r"<script[^>]*>.*?</script>", "", text, flags=re.DOTALL | re.I)
    text = re.sub(r"<br\s*/?>", "\n", text, flags=re.I)
    text = re.sub(r"</(?:p|div|tr|li)>", "\n", text, flags=re.I)
    text = re.sub(r"<[^>]+>", "", text)
    text = html.unescape(text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ── Size formatting ─────────────────────────────────────────────────


def fmt_size(n: int) -> str:
    if n < 1024:
        return f"{n} B"
    if n < 1024 * 1024:
        return f"{n / 1024:.1f} KB"
    return f"{n / (1024 * 1024):.1f} MB"


# ── Markdown report ─────────────────────────────────────────────────

PARSEABLE_EXTS = {".pdf", ".xlsx", ".xls", ".pptx", ".txt", ".csv", ".log", ".json", ".xml", ".html", ".htm", ".md"}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg"}


def build_report(filepath: str) -> str:
    headers, body_plain, body_html, attachments = parse_eml(filepath)

    out = []
    out.append("# 이메일 분석서\n")
    out.append(f"**원본 파일**: `{os.path.basename(filepath)}`\n")

    # ── 메일 정보
    out.append("## 메일 정보\n")
    out.append("| 항목 | 내용 |")
    out.append("|------|------|")
    out.append(f"| **제목** | {headers['subject']} |")
    out.append(f"| **보낸 사람** | {headers['from']} |")
    out.append(f"| **받는 사람** | {headers['to']} |")
    if headers["cc"]:
        out.append(f"| **참조** | {headers['cc']} |")
    out.append(f"| **날짜** | {headers['date']} |")
    out.append(f"| **첨부파일** | {len(attachments)}개 |")
    out.append("")

    # ── 본문
    out.append("## 본문 내용\n")
    body = body_plain
    if not body and body_html:
        body = strip_html(body_html)
    out.append(body.strip() if body else "_(본문 없음)_")
    out.append("")

    # ── 첨부파일
    if attachments:
        out.append("## 첨부파일 분석\n")
        out.append("| # | 파일명 | 형식 | 크기 |")
        out.append("|---|--------|------|------|")
        for i, a in enumerate(attachments, 1):
            out.append(f"| {i} | {a['filename']} | {a['content_type']} | {fmt_size(a['size'])} |")
        out.append("")

        for i, a in enumerate(attachments, 1):
            ext = Path(a["filename"]).suffix.lower()
            out.append(f"### 첨부 {i}: {a['filename']}\n")

            if ext == ".pdf":
                out.append(extract_pdf(a["data"]))
            elif ext in (".xlsx", ".xls"):
                out.append(extract_xlsx(a["data"]))
            elif ext == ".pptx":
                out.append(extract_pptx(a["data"]))
            elif ext == ".ppt":
                out.append("_(.ppt 레거시 형식 미지원, .pptx만 지원)_")
            elif ext in (".txt", ".csv", ".log", ".json", ".xml", ".html", ".htm", ".md"):
                out.append(f"```\n{extract_text_file(a['data'])}\n```")
            elif ext in IMAGE_EXTS:
                out.append(f"_(이미지 파일 - {fmt_size(a['size'])})_")
            else:
                out.append(f"_(지원하지 않는 형식: {ext}, {fmt_size(a['size'])})_")
            out.append("")

    return "\n".join(out)


# ── Main ────────────────────────────────────────────────────────────

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python eml-analyzer.py <eml_file_path>", file=sys.stderr)
        sys.exit(1)

    path = sys.argv[1]
    if not os.path.isfile(path):
        print(f"파일을 찾을 수 없습니다: {path}", file=sys.stderr)
        sys.exit(1)

    print(build_report(path))
