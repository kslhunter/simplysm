"""PPTX handler: extract text, images, and OLE embedded objects."""

from _common import ensure_packages, ext_from_content_type

PACKAGES = {"python-pptx": "pptx"}


def _emu_to_inches(emu):
    if emu is None:
        return "?"
    return f"{emu / 914400:.1f}"


def _pos(shape):
    return f"(left={_emu_to_inches(shape.left)}\", top={_emu_to_inches(shape.top)}\")"


def extract(file_path):
    ensure_packages(PACKAGES)
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    text_parts = []
    images = []
    embedded = []

    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts.append(f"[Slide {slide_num}]")

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ext = ext_from_content_type(shape.image.content_type)
                images.append({
                    "data": shape.image.blob,
                    "ext": ext,
                    "context": f"Slide {slide_num} {_pos(shape)}",
                })

            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip().replace("\n", "\n       ")
                text_parts.append(f"[TXT] {_pos(shape)} {text}")

        # OLE embedded objects from slide relationships
        seen = set()
        for rel in slide.part.rels.values():
            reltype = rel.reltype or ""
            if "oleObject" in reltype or "package" in reltype:
                target_ref = getattr(rel, 'target_ref', '') or ''
                if target_ref in seen:
                    continue
                seen.add(target_ref)
                try:
                    blob = rel.target_part.blob
                    filename = target_ref.split("/")[-1] if "/" in target_ref else target_ref
                    if not filename:
                        filename = f"embedded_slide{slide_num}_{len(embedded) + 1}.bin"
                    embedded.append({"filename": filename, "data": blob})
                except Exception:
                    pass

        text_parts.append("")

    return {
        "text": "\n".join(text_parts),
        "images": images,
        "embedded": embedded,
        "metadata": {},
    }
