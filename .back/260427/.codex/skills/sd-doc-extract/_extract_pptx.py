"""PPTX handler: render slides to PNG via PowerPoint COM, extract text and OLE embedded.

Individual image/shape extraction is intentionally omitted — slide screenshots
contain all visuals including overlay shapes (boxes, arrows, annotations) that
lose their spatial relationship when decomposed. Requires Windows + Microsoft
PowerPoint installed.
"""

import tempfile
from pathlib import Path

from _common import ensure_packages

PACKAGES = {"pywin32": "win32com.client", "python-pptx": "pptx"}


def _emu_to_inches(emu):
    if emu is None:
        return "?"
    return f"{emu / 914400:.1f}"


def _pos(shape):
    return f"(left={_emu_to_inches(shape.left)}\", top={_emu_to_inches(shape.top)}\")"


def _extract_shapes(shapes, text_parts):
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            _extract_shapes(shape.shapes, text_parts)
        elif shape.has_table:
            tbl = shape.table
            text_parts.append(f"[TABLE] {_pos(shape)}")
            for r_idx, row in enumerate(tbl.rows):
                cells = [
                    cell.text.strip().replace("\\", "\\\\").replace("|", "\\|")
                    .replace("\r\n", "<br>").replace("\n", "<br>").replace("\r", "<br>")
                    for cell in row.cells
                ]
                text_parts.append("| " + " | ".join(cells) + " |")
                if r_idx == 0:
                    text_parts.append("|" + "|".join(["---"] * len(cells)) + "|")
        elif hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip().replace("\n", "\n       ")
            text_parts.append(f"[TXT] {_pos(shape)} {text}")


def _render_slides_via_com(file_path: str, tmp_dir: Path, slide_count: int,
                           width: int, height: int) -> list[bytes]:
    import win32com.client
    import pythoncom

    pythoncom.CoInitialize()
    try:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        try:
            try:
                app.DisplayAlerts = 0
            except Exception:
                pass
            abs_path = str(Path(file_path).resolve())
            prs = app.Presentations.Open(abs_path, ReadOnly=True, Untitled=False,
                                         WithWindow=False)
            try:
                results = []
                for i in range(1, slide_count + 1):
                    tmp_path = tmp_dir / f"__tmp_slide_{i}.png"
                    prs.Slides(i).Export(str(tmp_path), "PNG", width, height)
                    results.append(tmp_path.read_bytes())
                    tmp_path.unlink()
                return results
            finally:
                prs.Close()
        finally:
            app.Quit()
    finally:
        pythoncom.CoUninitialize()


def extract(file_path):
    ensure_packages(PACKAGES)
    from pptx import Presentation

    prs = Presentation(file_path)
    slide_count = len(prs.slides)

    target_width = 1920
    if prs.slide_width and prs.slide_height:
        target_height = int(target_width * prs.slide_height / prs.slide_width)
    else:
        target_height = 1080

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            slide_pngs = _render_slides_via_com(
                file_path, Path(tmpdir), slide_count, target_width, target_height
            )
    except Exception as e:
        raise RuntimeError(
            f"PowerPoint COM rendering failed: {e}. "
            "This extractor requires Windows with Microsoft PowerPoint installed."
        ) from e

    text_parts = []
    slide_images = []
    embedded = []
    emb_idx = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts.append(f"[Slide {slide_num}]")

        slide_images.append({
            "filename": f"slide_{slide_num:03d}.png",
            "data": slide_pngs[slide_num - 1],
        })
        text_parts.append(f"[SLIDE:{slide_num}]")

        _extract_shapes(slide.shapes, text_parts)

        # Speaker notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip() if notes_frame else ""
            if notes_text:
                notes_text = notes_text.replace("\n", "\n       ")
                text_parts.append(f"[Notes] {notes_text}")

        seen = set()
        for rel in slide.part.rels.values():
            reltype = rel.reltype or ""
            if "oleObject" in reltype or "package" in reltype:
                target_ref = getattr(rel, 'target_ref', '') or ''
                if target_ref in seen:
                    continue
                seen.add(target_ref)
                try:
                    blob = rel.target_part.blob
                    filename = target_ref.split("/")[-1] if "/" in target_ref else target_ref
                    if not filename:
                        filename = f"embedded_slide{slide_num}_{len(embedded) + 1}.bin"
                    emb_idx += 1
                    embedded.append({"filename": filename, "data": blob})
                    text_parts.append(f"[EMB:{emb_idx}]")
                except Exception:
                    pass

        text_parts.append("")

    return {
        "text": "\n".join(text_parts),
        "images": [],
        "embedded": embedded,
        "metadata": {},
        "slide_images": slide_images,
    }
